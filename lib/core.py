# NOTE!!NOTE!!!NOTE!!NOTE!!!NOTE!!NOTE!!!NOTE!!NOTE!!!
# THE WORD "CHAPTER" IN THE CODE DOES NOT MEAN
# IT'S THE REAL CHAPTER OF THE EBOOK SINCE NO STANDARDS
# ARE DEFINING A CHAPTER ON .EPUB FORMAT. THE WORD "BLOCK"
# IS USED TO PRINT IT OUT TO THE TERMINAL, AND "CHAPTER" TO THE CODE
# WHICH IS LESS GENERIC FOR THE DEVELOPERS

import argparse, asyncio, csv, fnmatch, sqlite3, hashlib, io, json, math, os, pytesseract, gc
import random, shutil, subprocess, sys, tempfile, threading, time, uvicorn, copy
import traceback, socket, unicodedata, urllib.request, uuid, zipfile, fitz, multiprocessing
import ebooklib, psutil, requests, stanza, importlib, queue
import regex as re, gradio as gr

from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import Any, Generator, Dict
from PIL import Image, ImageSequence
from tqdm import tqdm
from bs4 import BeautifulSoup, NavigableString, Tag
from collections import Counter
from collections.abc import Mapping, MutableMapping
from datetime import datetime
from ebooklib import epub
from ebooklib.epub import EpubBook
from ebooklib.epub import EpubHtml
from glob import glob
from iso639 import Lang
from markdown import markdown
from multiprocessing import Pool, cpu_count
from multiprocessing import Manager, Event
from multiprocessing.managers import DictProxy, ListProxy, SyncManager
from stanza.pipeline.core import Pipeline, DownloadMethod
from num2words2 import num2words
from pathlib import Path
from PIL import Image
from pydub import AudioSegment
from pydub.utils import mediainfo
from queue import Queue, Empty
from types import MappingProxyType
from langdetect import detect
from unidecode import unidecode
from phonemizer import phonemize

from lib.classes.subprocess_pipe import SubprocessPipe
from lib.classes.vram_detector import VRAMDetector
from lib.classes.voice_extractor import VoiceExtractor
from lib.classes.tts_manager import TTSManager
#from lib.classes.redirect_console import RedirectConsole
#from lib.classes.argos_translator import ArgosTranslator
from lib.classes.tts_engines.common.audio import get_audiolist_duration, get_audio_duration
from lib.classes.tts_engines.common.utils import build_vtt_file

from lib import *

#import logging
#logging.basicConfig(
#    level=logging.INFO, # DEBUG for more verbosity
#    format="%(asctime)s [%(levelname)s] %(message)s"
#)

context = None
context_tracker = None
active_sessions = None
progress_bar = None

status_tags = {
    "OVERRIDE": "override",
    "DELETION": "deletion",
    "READY": "ready",
    "EDIT": "edit",
    "SKIP": "skip",
    "SWITCH": "switch",
    "CONVERTING": "converting",
    "END": "end",
    "DISCONNECTED": "disconnected"
}

ebook_modes = {
    "SINGLE": "single",
    "DIRECTORY": "directory",
    "TEXT": "text"
}

save_session_keys_except = [
    'blocks_orig',
    'blocks_saved',
    'blocks_current'
]

file_prefixes = {
    "clone": "__",
    "saved": "__saved_",
    'current': "__current_"
}

########### Classes

class DependencyError(Exception):
    def __init__(self, message:str|None):
        super().__init__(message)
        print(message)
        # Automatically handle the exception when it's raised
        self.handle_exception()

    def handle_exception(self)->None:
        # Print the full traceback of the exception
        traceback.print_exc()      
        # Print the exception message
        error = f'Caught DependencyError: {self}'
        print(error)

class SessionTracker:
    def __init__(self):
        self.lock = threading.Lock()
        #self.blocks_autosave = AppAutosave()
        #self.blocks_autosave.start()

    def start_session(self, session_id:str)->bool:
        with self.lock:
            if session_id not in context.sessions:
                return False
            session = context.get_session(session_id)
            session['status'] = status_tags['READY']
            return True

    def end_session(self, session_id:str, socket_hash:str)->None:
        #self.blocks_autosave.unregister(session_id)
        active_sessions.discard(socket_hash)
        with self.lock:
            context.sessions.pop(session_id, None)

class SessionContext:
    def __init__(self):
        self.manager:Manager = Manager()
        self.sessions:DictProxy[str, DictProxy[str, Any]] = self.manager.dict()
        self.cancellation_events = {}

    def _recursive_proxy(self, data:Any, manager:SyncManager|None)->Any:
        if manager is None:
            manager = self.manager
        if isinstance(data, dict):
            proxy_dict = manager.dict()
            for key, value in data.items():
                proxy_dict[key] = self._recursive_proxy(value, manager)
            return proxy_dict
        elif isinstance(data, list):
            proxy_list = manager.list()
            for item in data:
                proxy_list.append(self._recursive_proxy(item, manager))
            return proxy_list
        elif isinstance(data, (str, int, float, bool, type(None))):
            return data
        else:
            error = f'Unsupported data type: {type(data)}'
            print(error)
            return None

    def set_session(self, session_id:str)->Any:
        self.sessions[session_id] = self._recursive_proxy({
            ####### Global settings
            "id": session_id,
            "script_mode": NATIVE,
            "tab_id": None,
            "socket_hash": None,
            "session_dir": None,
            "is_gui_process": False,
            "free_vram_gb": 0,
            "status": None,
            "ticker": 0,
            "cancellation_requested": False,
            "ebook_mode": ebook_modes['SINGLE'],
            "blocks_preview": False,
            "device": default_device,
            "tts_engine": default_tts_engine,
            "fine_tuned": default_fine_tuned,
            "model_cache": None,
            "model_zs_cache": None,
            "stanza_cache": None,
            "system": None,
            "client": None,
            "language": default_language_code,
            "language_iso1": None,
            "voice": None,
            "voice_dir": None,
            "voice_map": {},
            "ebook_selected": None,
            "custom_model": None,
            "custom_model_dir": None,
            "output_dir": None,
            "output_format": default_output_format,
            "output_channel": default_output_channel,
            "output_split": default_output_split,
            "output_split_hours": default_output_split_hours,
            ####### Xtts settings
            "xtts_temperature": default_engine_settings[TTS_ENGINES['XTTSv2']]['temperature'],
            #"xtts_codec_temperature": default_engine_settings[TTS_ENGINES['XTTSv2']]['codec_temperature'],
            "xtts_length_penalty": default_engine_settings[TTS_ENGINES['XTTSv2']]['length_penalty'],
            "xtts_num_beams": default_engine_settings[TTS_ENGINES['XTTSv2']]['num_beams'],
            "xtts_repetition_penalty": default_engine_settings[TTS_ENGINES['XTTSv2']]['repetition_penalty'],
            #"xtts_cvvp_weight": default_engine_settings[TTS_ENGINES['XTTSv2']]['cvvp_weight'],
            "xtts_top_k": default_engine_settings[TTS_ENGINES['XTTSv2']]['top_k'],
            "xtts_top_p": default_engine_settings[TTS_ENGINES['XTTSv2']]['top_p'],
            "xtts_speed": default_engine_settings[TTS_ENGINES['XTTSv2']]['speed'],
            #"xtts_gpt_cond_len": default_engine_settings[TTS_ENGINES['XTTSv2']]['gpt_cond_len'],
            #"xtts_gpt_batch_size": default_engine_settings[TTS_ENGINES['XTTSv2']]['gpt_batch_size'],
            "xtts_enable_text_splitting": default_engine_settings[TTS_ENGINES['XTTSv2']]['enable_text_splitting'],
            ####### Bark settings
            "bark_text_temp": default_engine_settings[TTS_ENGINES['BARK']]['text_temp'],
            "bark_waveform_temp": default_engine_settings[TTS_ENGINES['BARK']]['waveform_temp'],
            ####### Audiobook editor
            "audiobook": None,
            "audiobooks_dir": None,
            ####### Ebook conversion
            "ebook": None,
            "ebook_src": None,
            "ebook_src_notextarea": None,
            "ebook_list": None,
            "ebook_textarea": None,
            "audiobook_overridden": None,
            "process_dir": None,
            "chapters_dir": None,
            "sentences_dir": None,
            "epub_path": None,
            "final_name": None,
            "filename_noext": None,
            "cover": None,
            "blocks_orig": {},
            "blocks_saved": {},
            "blocks_current": {},
            "blocks_orig_json": None,
            "blocks_saved_json": None,
            "blocks_current_db": None,
            "duration": 0,
            "playback_time": 0,
            "playback_volume": 0,
            "metadata": {
                "title": None, 
                "creator": None,
                "contributor": None,
                "language": None,
                "identifier": None,
                "publisher": None,
                "date": None,
                "description": None,
                "subject": None,
                "rights": None,
                "format": None,
                "type": None,
                "coverage": None,
                "relation": None,
                "Source": None,
                "Modified": None,
            }
        }, manager=self.manager)
        return self.sessions[session_id]

    def get_session(self, session_id:str)->Any:
        if session_id in self.sessions:
            return self.sessions[session_id]
        return {}

    def find_id_by_hash(self, socket_hash:str)->str|None:
        for session_id, session in list(self.sessions.items()):
            if socket_hash in session:
                return session_id
        return None
        
class JSONDictProxyEncoder(json.JSONEncoder):
    def default(self, o:Any)->Any:
        if isinstance(o, DictProxy):
            return dict(o)
        elif isinstance(o, ListProxy):
            return list(o)
        return super().default(o)

"""
class AppAutosave:
    def __init__(self, interval:float=15.0):
        self._interval = interval
        self._sessions: set[str] = set()
        self._lock = threading.Lock()
        self._started = False

    def start(self)->None:
        if self._started:
            return
        self._started = True
        t = threading.Thread(target=self._timer, daemon=True)
        t.start()

    def register(self, session_id:str)->None:
        with self._lock:
            self._sessions.add(session_id)

    def unregister(self, session_id:str)->None:
        with self._lock:
            self._sessions.discard(session_id)

    def _timer(self)->None:
        while True:
            time.sleep(self._interval)
            with self._lock:
                session_ids = set(self._sessions)
            for session_id in session_ids:
                try:
                    session = context.get_session(session_id)
                    if not session or not session.get('id', False):
                        with self._lock:
                            self._sessions.discard(session_id)
                        continue
                except Exception as e:
                    logger.error(f'AppAutosave._timer({session_id}): {e}!')
"""
        
############# End classes

def prepare_dirs(session_id:str)->bool:
    try:
        session = context.get_session(session_id)
        if session and session.get('id', False):
            resume = False
            os.makedirs(os.path.join(models_dir,'tts'), exist_ok=True)
            os.makedirs(session['session_dir'], exist_ok=True)
            os.makedirs(session['process_dir'], exist_ok=True)
            os.makedirs(session['custom_model_dir'], exist_ok=True)
            os.makedirs(session['voice_dir'], exist_ok=True)
            os.makedirs(session['audiobooks_dir'], exist_ok=True)
            os.makedirs(session['chapters_dir'], exist_ok=True)
            os.makedirs(session['sentences_dir'], exist_ok=True)
            return True
    except Exception as e:
        DependencyError(e)
        return False

def check_programs(prog_name:str, command:str, options:str)->bool:
    try:
        subprocess.run(
            [command, options],
            stdout=subprocess.PIPE, 
            stderr=subprocess.PIPE,
            check=True,
            text=True,
            encoding='utf-8'
        )
        return True
    except FileNotFoundError:
        e = f'''********** Error: {prog_name} is not installed! if your OS calibre package version 
        is not compatible you still can run ebook2audiobook.sh (linux/mac) or ebook2audiobook.cmd (windows) **********'''
        DependencyError(e)
    except subprocess.CalledProcessError:
        e = f'Error: There was an issue running {prog_name}.'
        DependencyError(e)
    return False

def analyze_uploaded_file(zip_path:str, required_files:list[str])->bool:
    try:
        if not os.path.exists(zip_path):
            error = f'The file does not exist: {os.path.basename(zip_path)}'
            print(error)
            return False
        files_in_zip = {}
        empty_files = set()
        with zipfile.ZipFile(zip_path, 'r') as zf:
            for file_info in zf.infolist():
                file_name = file_info.filename
                if file_info.is_dir():
                    continue
                base_name = os.path.basename(file_name)
                files_in_zip[base_name] = file_info.file_size
                if file_info.file_size == 0:
                    empty_files.add(base_name)
        required_files = [file for file in required_files]
        missing_files = [f for f in required_files if f not in files_in_zip]
        required_empty_files = [f for f in required_files if f in empty_files]
        if missing_files:
            msg = f'Missing required files: {missing_files}'
            print(msg)
        if required_empty_files:
            msg = f'Required files with 0 KB: {required_empty_files}'
            print(msg)
        return not missing_files and not required_empty_files
    except zipfile.BadZipFile:
        error = 'The file is not a valid ZIP archive.'
        print(error)
        return False
    except Exception as e:
        error = f'An error occurred: {e}'
        print(error)
        return False

def extract_custom_model(session_id)->str|None:
    session = context.get_session(session_id)
    if session and session.get('id', False):
        file_src = session['custom_model']
        required_files = default_engine_settings[session['tts_engine']]['files']
        model_path = None
        model_name = re.sub('.zip', '', os.path.basename(file_src), flags=re.IGNORECASE)
        model_name = get_sanitized(model_name)
        try:
            with zipfile.ZipFile(file_src, 'r') as zip_ref:
                files = zip_ref.namelist()
                files_length = len(files)
                tts_dir = session['tts_engine']
                model_path = os.path.join(session['custom_model_dir'], tts_dir, model_name)
                os.makedirs(model_path, exist_ok=True)
                msg = f'Extracting files to {model_path}…'
                with tqdm(total=files_length, unit='files') as t:
                    for f in files:
                        base_f = os.path.basename(f)
                        if base_f in required_files:
                            out_path = os.path.join(model_path, base_f)
                            with zip_ref.open(f) as src, open(out_path, 'wb') as dst:
                                shutil.copyfileobj(src, dst)
                        t.update(1)
                        if session['is_gui_process']:
                            progress_bar((t.n + 1) / files_length, desc=msg)
            if model_path is not None:
                if session['tts_engine'] in tts_engines_with_inner_speaker:
                    if os.path.exists(file_src):
                        os.remove(file_src)
                    return model_path
                msg = f'Normalizing ref.wav…'
                print(msg)
                voice_ref = os.path.join(model_path, 'ref.wav')
                voice_name = model_name
                final_voice_file = os.path.join(model_path, f'{voice_name}.wav')
                extractor = VoiceExtractor(session, voice_ref, voice_name, final_voice_file)
                status, msg = extractor.extract_voice()
                if status:
                    session['voice'] = final_voice_file
                    if os.path.exists(file_src):
                        os.remove(file_src)
                    if os.path.exists(voice_ref):
                        os.remove(voice_ref)
                    return model_path
                else:
                    error = f'extract_custom_model() VoiceExtractor.extract_voice() error! {msg}'
                    print(error)
            else:
                error = f'An error occurred     when unzip {file_src}'
                print(error)
        except asyncio.exceptions.CancelledError as e:
            DependencyError(e)
            error = f'extract_custom_model asyncio.exceptions.CancelledError: {e}'
            print(error)
        except Exception as e:
            DependencyError(e)
            error = f'extract_custom_model Exception: {e}'
            print(error)
        if session['is_gui_process']:
            if os.path.exists(file_src):
                os.remove(file_src)
        session['custom_model'] = None
    return None
        
def hash_proxy_dict(proxy_dict:Any)->str:
    try:
        data = {k: v for k, v in dict(proxy_dict).items() if k not in save_session_keys_except}
    except Exception:
        data = {}
    data_str = json.dumps(data, default=str, sort_keys=True)
    return hashlib.md5(data_str.encode('utf-8')).hexdigest()

def compare_checksums(session_id:str)->tuple[bool, str|None]:
    try:
        session = context.get_session(session_id)
        if session and session.get('id', False):
            hash_algorithm:str = 'sha256'
            checksum_path = os.path.join(session['process_dir'], 'checksum')
            hash_func = hashlib.new(hash_algorithm)
            with open(session['ebook'], 'rb') as f:
                while chunk := f.read(8192):
                    hash_func.update(chunk)
            new_checksum = hash_func.hexdigest()
            if not os.path.exists(checksum_path):
                with open(checksum_path, 'w', encoding='utf-8') as f:
                    f.write(new_checksum)
                return False, None
            else:
                with open(checksum_path, 'r', encoding='utf-8') as f:
                    saved_checksum = f.read().strip()
                if saved_checksum == new_checksum:
                    return True, None
                else:
                    with open(checksum_path, 'w', encoding='utf-8') as f:
                        f.write(new_checksum)
                        return False, None
        error = f'compare_checksums() error: session does not exist'
        return False, error
    except Exception as e:
        error = f'compare_checksums() error: {e}'
        return False, error

def compare_dict_keys(d1, d2):
    if not isinstance(d1, Mapping) or not isinstance(d2, Mapping):
        return d1 == d2
    d1_keys = set(d1.keys())
    d2_keys = set(d2.keys())
    missing_in_d2 = d1_keys - d2_keys
    missing_in_d1 = d2_keys - d1_keys
    if missing_in_d2 or missing_in_d1:
        return {
            "missing_in_d2": missing_in_d2,
            "missing_in_d1": missing_in_d1,
        }
    for key in d1_keys.intersection(d2_keys):
        nested_result = compare_keys(d1[key], d2[key])
        if nested_result:
            return {key: nested_result}
    return None

def ocr2xhtml(img: Image.Image, lang:str)->tuple[str|bool, str|None]:
    try:
        debug = True
        try:
            data = pytesseract.image_to_data(img, lang=lang, output_type=pytesseract.Output.DATAFRAME)
            # Handle silent OCR failures (empty or None result)
            if data is None or data.empty:
                error = f'Tesseract returned empty OCR data for language "{lang}".'
                return False, error
        except (pytesseract.TesseractError, Exception) as e:
            print(f'The OCR {lang} trained model must be downloaded.')
            try:
                tessdata_dir = os.environ['TESSDATA_PREFIX']
                os.makedirs(tessdata_dir, exist_ok=True)
                url = f'https://github.com/tesseract-ocr/tessdata_best/raw/main/{lang}.traineddata'
                dest_path = os.path.join(tessdata_dir, f'{lang}.traineddata')
                msg = f'Downloading {lang}.traineddata into {tessdata_dir}…'
                print(msg)
                response = requests.get(url, timeout=15)
                if response.status_code == 200:
                    with open(dest_path, 'wb') as f:
                        f.write(response.content)
                    msg = f'Downloaded and installed {lang}.traineddata successfully.'
                    print(msg)
                    data = pytesseract.image_to_data(img, lang=lang, output_type=pytesseract.Output.DATAFRAME)
                    if data is None or data.empty:
                        error = f'Tesseract returned empty OCR data even after downloading {lang}.traineddata.'
                        return False, error
                else:
                    error = f'Failed to download traineddata for {lang} (HTTP {response.status_code})'
                    return False, error
            except Exception as e:
                error = f'Automatic download failed: {e}'
                return False, error
        data = data.dropna(subset=['text'])
        lines = []
        last_block = None
        for _, row in data.iterrows():
            text = row['text'].strip()
            if not text:
                continue
            block = row['block_num']
            if last_block is not None and block != last_block:
                lines.append('')  # blank line between blocks
            lines.append(text)
            last_block = block
        joined = '\n'.join(lines)
        raw_lines = [l.strip() for l in joined.split('\n')]
        # Normalize line breaks
        merged_lines = []
        buffer = ''
        for i, line in enumerate(raw_lines):
            if not line:
                if buffer:
                    merged_lines.append(buffer.strip())
                    buffer = ''
                continue
            if buffer and not buffer.endswith(('.', '?', '!', ':')) and not line[0].isupper():
                buffer += ' ' + line
            else:
                if buffer:
                    merged_lines.append(buffer.strip())
                buffer = line
        if buffer:
            merged_lines.append(buffer.strip())
        # Detect heading-like lines
        xhtml_parts = []
        debug_dump = []
        for i, p in enumerate(merged_lines):
            is_heading = False
            if p.isupper() and len(p.split()) <= 8:
                is_heading = True
            elif len(p.split()) <= 5 and p.istitle():
                is_heading = True
            elif (i == 0 or (i > 0 and merged_lines[i-1] == '')) and len(p.split()) <= 10:
                is_heading = True
            if is_heading:
                xhtml_parts.append(f'<h2>{p}</h2>')
                debug_dump.append(f'[H2] {p}')
            else:
                xhtml_parts.append(f'<p>{p}</p>')
                debug_dump.append(f'[P ] {p}')
        if debug:
            print('=== OCR DEBUG OUTPUT ===')
            for line in debug_dump:
                print(line)
            print('========================')
        return '\n'.join(xhtml_parts), None
    except Exception as e:
        DependencyError(e)
        error = f'ocr2xhtml error: {e}'
        return False, error

def create_db_blocks(db_path:str)->None:
    os.makedirs(os.path.dirname(db_path) or '.', exist_ok=True)
    with sqlite3.connect(db_path) as conn:
        conn.execute('PRAGMA journal_mode=WAL')
        conn.execute('PRAGMA synchronous=NORMAL')
        conn.execute('PRAGMA foreign_keys=ON')
        conn.executescript('''
            CREATE TABLE IF NOT EXISTS stamp (
                id INTEGER PRIMARY KEY CHECK (id = 1),
                page INTEGER,
                block_resume INTEGER,
                sentence_resume INTEGER,
                voice TEXT,
                tts_engine TEXT,
                fine_tuned TEXT
            );
            CREATE TABLE IF NOT EXISTS blocks (
                id TEXT PRIMARY KEY,
                idx INTEGER NOT NULL,
                expand INTEGER NOT NULL,
                keep INTEGER NOT NULL,
                text TEXT NOT NULL,
                voice TEXT,
                tts_engine TEXT,
                fine_tuned TEXT
            );
            CREATE TABLE IF NOT EXISTS sentences (
                block_id TEXT NOT NULL,
                idx INTEGER NOT NULL,
                text TEXT NOT NULL,
                PRIMARY KEY (block_id, idx),
                FOREIGN KEY (block_id) REFERENCES blocks(id) ON DELETE CASCADE
            );
            CREATE INDEX IF NOT EXISTS idx_blocks_idx ON blocks(idx);
            INSERT OR IGNORE INTO stamp (id, page, block_resume, sentence_resume, voice, tts_engine, fine_tuned)
            VALUES (1, 0, 0, 0, NULL, NULL, NULL);
        ''')

def load_db_blocks(db_path:str)->dict:
    try:
        if not os.path.exists(db_path):
            return {}
        with sqlite3.connect(db_path) as conn:
            conn.execute('PRAGMA foreign_keys=ON')
            stamp_row = conn.execute(
                'SELECT page, block_resume, sentence_resume, voice, tts_engine, fine_tuned FROM stamp WHERE id=1'
            ).fetchone()
            if stamp_row is None:
                return {}
            page, block_resume, sentence_resume, voice, tts_engine, fine_tuned = stamp_row
            sentences_by_block = {}
            for block_id, text in conn.execute('SELECT block_id, text FROM sentences ORDER BY block_id, idx'):
                sentences_by_block.setdefault(block_id, []).append(text)
            blocks = []
            for row in conn.execute('SELECT id, expand, keep, text, voice, tts_engine, fine_tuned FROM blocks ORDER BY idx'):
                bid, expand, keep, text, b_voice, b_tts_engine, b_fine_tuned = row
                blocks.append({
                    'id': bid,
                    'expand': bool(expand),
                    'keep': bool(keep),
                    'text': text,
                    'voice': b_voice,
                    'tts_engine': b_tts_engine,
                    'fine_tuned': b_fine_tuned,
                    'sentences': sentences_by_block.get(bid, []),
                })
            return {
                'page': page,
                'block_resume': block_resume,
                'sentence_resume': sentence_resume,
                'voice': voice,
                'tts_engine': tts_engine,
                'fine_tuned': fine_tuned,
                'blocks': blocks,
            }
    except Exception as e:
        error = f'load_db_blocks() error: {e}'
        print(error)
        return {}

def save_db_stamp(session_id:str)->None:
    try:
        session = context.get_session(session_id)
        if not (session and session.get('id', False)):
            return
        data = session['blocks_current']
        if not data:
            return
        db_path = session['blocks_current_db']
        create_db_blocks(db_path)
        with sqlite3.connect(db_path) as conn:
            conn.execute(
                'UPDATE stamp SET page=?, block_resume=?, sentence_resume=?, voice=?, tts_engine=?, fine_tuned=? WHERE id=1',
                (
                    data.get('page'),
                    data.get('block_resume'),
                    data.get('sentence_resume'),
                    data.get('voice'),
                    data.get('tts_engine'),
                    data.get('fine_tuned'),
                )
            )
            conn.commit()
    except Exception as e:
        error = f'save_db_stamp() error: {e}'
        print(error)

def save_db_blocks(session_id:str)->None:
    try:
        session = context.get_session(session_id)
        if not (session and session.get('id', False)):
            return
        data = session['blocks_current']
        if not data:
            return
        db_path = session['blocks_current_db']
        create_db_blocks(db_path)
        with sqlite3.connect(db_path) as conn:
            conn.execute('PRAGMA foreign_keys=ON')
            conn.execute(
                'UPDATE stamp SET page=?, block_resume=?, sentence_resume=?, voice=?, tts_engine=?, fine_tuned=? WHERE id=1',
                (
                    data.get('page'),
                    data.get('block_resume'),
                    data.get('sentence_resume'),
                    data.get('voice'),
                    data.get('tts_engine'),
                    data.get('fine_tuned'),
                )
            )
            new_blocks = data.get('blocks', [])
            new_ids = {b['id'] for b in new_blocks}
            existing_ids = {row[0] for row in conn.execute('SELECT id FROM blocks')}
            removed = existing_ids - new_ids
            if removed:
                conn.executemany('DELETE FROM blocks WHERE id=?', [(rid,) for rid in removed])
            for idx, block in enumerate(new_blocks):
                block_id = block['id']
                conn.execute(
                    'INSERT INTO blocks (id, idx, expand, keep, text, voice, tts_engine, fine_tuned) '
                    'VALUES (?, ?, ?, ?, ?, ?, ?, ?) '
                    'ON CONFLICT(id) DO UPDATE SET '
                    'idx=excluded.idx, expand=excluded.expand, keep=excluded.keep, text=excluded.text, '
                    'voice=excluded.voice, tts_engine=excluded.tts_engine, fine_tuned=excluded.fine_tuned',
                    (
                        block_id,
                        idx,
                        1 if block.get('expand') else 0,
                        1 if block.get('keep') else 0,
                        block.get('text', ''),
                        block.get('voice'),
                        block.get('tts_engine'),
                        block.get('fine_tuned'),
                    )
                )
                conn.execute('DELETE FROM sentences WHERE block_id=?', (block_id,))
                sentences = block.get('sentences', [])
                if sentences:
                    conn.executemany(
                        'INSERT INTO sentences (block_id, idx, text) VALUES (?, ?, ?)',
                        [(block_id, i, s) for i, s in enumerate(sentences)]
                    )
            conn.commit()
    except Exception as e:
        error = f'save_db_blocks() error: {e}'
        print(error)

def load_json_blocks(filepath:str)->dict:
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception as e:
        error = f"load_json_blocks() error: {e}"
        print(error)
        return {}

def save_json_blocks(session_id:str, key:str)->None:
    try:
        session = context.get_session(session_id)
        if (session and session.get('id', False)):
            json_data = session[key]
            json_path = session[f'{key}_json']
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(json_data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f'save_json_blocks() error: {e}')

def sync_globals_to_blocks(session_id:str)->None:
    try:
        session = context.get_session(session_id)
        if not (session and session.get('id', False)):
            return
        blocks_current = session.get('blocks_current') or {}
        anchor_voice = blocks_current.get('voice')
        current_voice = session.get('voice')
        if anchor_voice == current_voice:
            return
        changed = False
        for block in blocks_current.get('blocks', []):
            if block.get('voice') == anchor_voice:
                block['voice'] = current_voice
                changed = True
        blocks_current['voice'] = current_voice
        session['blocks_current'] = blocks_current
    except Exception as e:
        exception_alert(session_id, f'sync_globals_to_blocks(): {e}')

def convert2epub(session_id:str)-> bool:
    session = context.get_session(session_id)
    if session and session.get('id', False):
        if session['cancellation_requested']:
            return False
        try:
            title = False
            author = False
            ebook_convert = shutil.which('ebook-convert')
            if not ebook_convert:
                error = 'ebook-convert utility is not installed or not found.'
                print(error)
                return False
            file_input = session['ebook']
            if os.path.getsize(file_input) == 0:
                error = f'Input file is empty: {file_input}'
                print(error)
                return False
            file_ext = os.path.splitext(file_input)[1].lower()
            if file_ext not in ebook_formats:
                error = f'Unsupported file format: {file_ext}'
                print(error)
                return False
            if file_ext == '.txt':
                with open(file_input, 'r', encoding='utf-8') as f:
                    text = f.read()
                text = text.replace('\r\n', '\n')
                text = re.sub(r'\n{2,}', f".{TTS_SML['pause']['static']}", text)
                with open(file_input, 'w', encoding='utf-8') as f:
                    f.write(text)
            elif file_ext == '.pdf':
                msg = 'File input is a PDF. flatten it in XHTML…'
                print(msg)
                doc = fitz.open(file_input)
                file_meta = doc.metadata
                filename_noext = os.path.splitext(os.path.basename(session['ebook']))[0]
                title = file_meta.get('title') or filename_noext
                author = file_meta.get('author') or False
                xhtml_pages = []
                for i, page in enumerate(doc):
                    has_text = page.get_text('text').strip()
                    if has_text:
                        try:
                            xhtml_content = page.get_text('xhtml').strip()
                        except Exception as e:
                            print(f'Error extracting text from page {i+1}: {e}')
                            xhtml_content = ''
                        error = None
                    else:
                        xhtml_content = ''
                        error = None
                    if not xhtml_content:
                        msg = f'The page {i+1} seems to be image-based. Using OCR…'
                        show_alert(session_id, {"type": "warning", "msg": msg})
                        pix = page.get_pixmap(dpi=300)
                        img = Image.open(io.BytesIO(pix.tobytes('png')))
                        xhtml_content, error = ocr2xhtml(img, session['language'])
                    if xhtml_content:
                        xhtml_pages.append(xhtml_content)
                    else:
                        show_alert(session_id, {"type": "warning", "msg": error})
                if xhtml_pages:
                    xhtml_body = '\n'.join(xhtml_pages)
                    xhtml_text = (
                        '<?xml version="1.0" encoding="utf-8"?>\n'
                        '<html xmlns="http://www.w3.org/1999/xhtml">\n'
                        '<head>\n'
                        f'<meta charset="utf-8"/>\n<title>{title}</title>\n'
                        '</head>\n'
                        '<body>\n'
                        f'{xhtml_body}\n'
                        '</body>\n'
                        '</html>\n'
                    )
                    file_input = os.path.join(session['process_dir'], f'{filename_noext}.xhtml')
                    with open(file_input, 'w', encoding='utf-8') as html_file:
                        html_file.write(xhtml_text)
                else:
                    return False
            elif file_ext == '.pptx':
                from html import escape as html_escape
                from pptx import Presentation as PptxPresentation
                filename_noext = os.path.splitext(os.path.basename(session['ebook']))[0]
                msg = f'File input is a presentation ({file_ext}). Extracting content…'
                print(msg)
                prs = PptxPresentation(file_input)
                title = prs.core_properties.title or filename_noext
                author = prs.core_properties.author or False
                xhtml_pages = []
                for i, slide in enumerate(prs.slides):
                    slide_texts = []
                    slide_images = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    slide_texts.append(text)
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_texts = [c.text.strip() for c in row.cells if c.text.strip()]
                                if row_texts:
                                    slide_texts.append(' | '.join(row_texts))
                        try:
                            slide_images.append(shape.image.blob)
                        except (AttributeError, ValueError):
                            pass
                    if slide_texts:
                        xhtml_content = '\n'.join(f'<p>{html_escape(t)}</p>' for t in slide_texts)
                    elif slide_images:
                        msg = f'Slide {i+1} seems to be image-based. Using OCR…'
                        show_alert(session_id, {"type": "warning", "msg": msg})
                        xhtml_parts = []
                        for blob in slide_images:
                            img = Image.open(io.BytesIO(blob))
                            xhtml_content, error = ocr2xhtml(img, session['language'])
                            if xhtml_content:
                                xhtml_parts.append(xhtml_content)
                            else:
                                show_alert(session_id, {"type": "warning", "msg": error})
                        xhtml_content = '\n'.join(xhtml_parts) if xhtml_parts else ''
                    else:
                        xhtml_content = ''
                    if xhtml_content:
                        xhtml_pages.append(xhtml_content)
                if xhtml_pages:
                    xhtml_body = '\n'.join(xhtml_pages)
                    xhtml_text = (
                        '<?xml version="1.0" encoding="utf-8"?>\n'
                        '<html xmlns="http://www.w3.org/1999/xhtml">\n'
                        '<head>\n'
                        f'<meta charset="utf-8"/>\n<title>{title}</title>\n'
                        '</head>\n'
                        '<body>\n'
                        f'{xhtml_body}\n'
                        '</body>\n'
                        '</html>\n'
                    )
                    file_input = os.path.join(session['process_dir'], f'{filename_noext}.xhtml')
                    with open(file_input, 'w', encoding='utf-8') as html_file:
                        html_file.write(xhtml_text)
                else:
                    return False
            elif file_ext == '.docx':
                from docx import Document as DocxDocument
                filename_noext = os.path.splitext(os.path.basename(session['ebook']))[0]
                docx_doc = DocxDocument(file_input)
                all_text = ''.join(p.text.strip() for p in docx_doc.paragraphs)
                if not all_text:
                    for table in docx_doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                all_text += cell.text.strip()
                                if all_text:
                                    break
                            if all_text:
                                break
                        if all_text:
                            break
                if not all_text:
                    msg = f'File input is a DOCX with no extractable text. Extracting images for OCR…'
                    print(msg)
                    title = docx_doc.core_properties.title or filename_noext
                    author = docx_doc.core_properties.author or False
                    xhtml_pages = []
                    for rel in docx_doc.part.rels.values():
                        if 'image' in rel.reltype:
                            try:
                                img = Image.open(io.BytesIO(rel.target_part.blob))
                                xhtml_content, error = ocr2xhtml(img, session['language'])
                                if xhtml_content:
                                    xhtml_pages.append(xhtml_content)
                                else:
                                    show_alert(session_id, {"type": "warning", "msg": error})
                            except Exception as e:
                                print(f'Error processing embedded image: {e}')
                    if xhtml_pages:
                        xhtml_body = '\n'.join(xhtml_pages)
                        xhtml_text = (
                            '<?xml version="1.0" encoding="utf-8"?>\n'
                            '<html xmlns="http://www.w3.org/1999/xhtml">\n'
                            '<head>\n'
                            f'<meta charset="utf-8"/>\n<title>{title}</title>\n'
                            '</head>\n'
                            '<body>\n'
                            f'{xhtml_body}\n'
                            '</body>\n'
                            '</html>\n'
                        )
                        file_input = os.path.join(session['process_dir'], f'{filename_noext}.xhtml')
                        with open(file_input, 'w', encoding='utf-8') as html_file:
                            html_file.write(xhtml_text)
                    else:
                        return False
            elif file_ext in ['.png', '.jpg', '.jpeg', '.tif', '.tiff', '.bmp']:
                filename_noext = os.path.splitext(os.path.basename(session['ebook']))[0]
                msg = f'File input is an image ({file_ext}). Running OCR…'
                print(msg)
                img = Image.open(file_input)
                xhtml_pages = []
                page_count = 0
                for i, frame in enumerate(ImageSequence.Iterator(img)):
                    page_count += 1
                    frame = frame.convert('RGB')
                    xhtml_content, error = ocr2xhtml(frame, session['language'])
                    if xhtml_content:
                        xhtml_pages.append(xhtml_content)
                    else:
                        show_alert(session_id, {"type": "warning", "msg": error})
                if xhtml_pages:
                    xhtml_body = '\n'.join(xhtml_pages)
                    xhtml_text = (
                        '<?xml version="1.0" encoding="utf-8"?>\n'
                        '<html xmlns="http://www.w3.org/1999/xhtml">\n'
                        '<head>\n'
                        f'<meta charset="utf-8"/>\n<title>{filename_noext}</title>\n'
                        '</head>\n'
                        '<body>\n'
                        f'{xhtml_body}\n'
                        '</body>\n'
                        '</html>\n'
                    )
                    file_input = os.path.join(session['process_dir'], f'{filename_noext}.xhtml')
                    with open(file_input, 'w', encoding='utf-8') as html_file:
                        html_file.write(xhtml_text)
                    print(f'OCR completed for {page_count} image page(s).')
                else:
                    return False
            msg = f"Running command: {ebook_convert} {file_input} {session['epub_path']}"
            print(msg)
            cmd = [
                    ebook_convert, file_input, session['epub_path'],
                    '--input-encoding=utf-8',
                    '--output-profile=generic_eink',
                    '--flow-size=0',
                    '--chapter-mark=pagebreak',
                    '--page-breaks-before',
                    "//*[name()='h1' or name()='h2' or name()='h3' or name()='h4' or name()='h5']",
                    '--disable-font-rescaling',
                    '--pretty-print',
                    '--smarten-punctuation',
                    '--verbose'
                ]
            if title:
                cmd += ['--title', title]
            if author:
                cmd += ['--authors', author]
            result = subprocess.run(
                cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding='utf-8'
            )
            print(result.stdout)
            return True
        except subprocess.CalledProcessError as e:
            DependencyError(e)
            error = f'convert2epub subprocess.CalledProcessError: {e.stderr}'
            print(error)
            return False
        except FileNotFoundError as e:
            DependencyError(e)
            error = f'convert2epub FileNotFoundError: {e}'
            print(error)
            return False
        except Exception as e:
            DependencyError(e)
            error = f'convert2epub error: {e}'
            print(error)
            return False

def get_ebook_title(epubBook:EpubBook,all_docs:list[Any])->str|None:
    # 1. Try metadata (official EPUB title)
    meta_title = epubBook.get_metadata('DC','title')
    if meta_title and meta_title[0][0].strip():
        return meta_title[0][0].strip()
    # 2. Try <title> in the head of the first XHTML document
    if all_docs:
        html = all_docs[0].get_content().decode('utf-8')
        soup = BeautifulSoup(html,'html.parser')
        title_tag = soup.select_one('head > title')
        if title_tag and title_tag.text.strip():
            return title_tag.text.strip()
        # 3. Try <img alt = '…'> if no visible <title>
        img = soup.find('img',alt = True)
        if img:
            alt = img['alt'].strip()
            if alt and 'cover' not in alt.lower():
                return alt
    return None

def get_cover(epubBook:EpubBook, session_id:str)->bool|str:
    try:
        session = context.get_session(session_id)
        if session and session.get('id', False):
            if session['cancellation_requested']:
                return False
            cover_image = None
            cover_path = os.path.join(session['process_dir'], session['filename_noext'] + '.jpg')
            for item in epubBook.get_items_of_type(ebooklib.ITEM_COVER):
                cover_image = item.get_content()
                break
            if not cover_image:
                for item in epubBook.get_items_of_type(ebooklib.ITEM_IMAGE):
                    if 'cover' in item.file_name.lower() or 'cover' in item.get_id().lower():
                        cover_image = item.get_content()
                        break
            if cover_image:
                # Open the image from bytes
                image = Image.open(io.BytesIO(cover_image))
                # Convert to RGB if needed (JPEG doesn't support alpha)
                if image.mode in ('RGBA', 'P'):
                    image = image.convert('RGB')
                image.save(cover_path, format = 'JPEG')
                return cover_path
            return True
    except Exception as e:
        DependencyError(e)
        return False

def get_blocks(session_id:str, epubBook:EpubBook)->list:
    try:
        msg = r'''
*******************************************************************************
NOTE:
The warning "Character xx not found in the vocabulary."
MEANS THE MODEL CANNOT INTERPRET THE CHARACTER AND WILL MAYBE GENERATE
(AS WELL AS WRONG PUNCTUATION POSITION) AN HALLUCINATION. THE BEST SOLUTION IS
TO MANUALLY REMOVE ALL UNRECOGNIZED CHARS AND WRONG PUNCTUATIONS FROM YOUR EBOOK
AND RESTART THE CONVERSION. TO IMPROVE THIS MODEL, IT NEEDS TO ADD THIS CHARACTER
INTO A NEW TRAINING MODEL. YOU CAN IMPROVE IT OR ASK TO A TRAINING MODEL EXPERT.
*******************************************************************************
        '''
        print(msg)
        session = context.get_session(session_id)
        if session and session.get('id', False):
            if session['cancellation_requested']:
                return []
            # Step 1: Extract TOC (Table of Contents)
            try:
                toc = epubBook.toc
                toc_list = [
                        nt for item in toc if hasattr(item, 'title')
                        if (nt := normalize_text(str(item.title), session['language'], session['language_iso1'], session['tts_engine'])) is not None
                ]
            except Exception as toc_error:
                error = f'Error extracting Table of Content: {toc_error}'
                show_alert(session_id, {"type": "warning", "msg": error})
            # Get spine item IDs
            spine_ids = [item[0] for item in epubBook.spine]
            # Filter only spine documents (i.e., reading order)
            all_docs = [
                item for item in epubBook.get_items_of_type(ebooklib.ITEM_DOCUMENT)
                if item.id in spine_ids
            ]
            if not all_docs:
                error = 'No document body found!'
                print(error)
                return []
            title = get_ebook_title(epubBook, all_docs)
            blocks = []
            stanza_nlp = False
            if session['language'] in year_to_decades_languages:
                try:
                    stanza_model = f"stanza-{session['language_iso1']}"
                    stanza_nlp = loaded_tts.get(stanza_model, False)
                    if stanza_nlp:
                        msg = f"NLP model {stanza_model} loaded."
                        print(msg)
                    else:
                        use_gpu = True if (
                            (session['device'] == devices['CUDA']['proc'] and devices['CUDA']['found']) or
                            (session['device'] == devices['ROCM']['proc'] and devices['ROCM']['found']) or
                            (session['device'] == devices['XPU']['proc'] and devices['XPU']['found']) or
                            (session['device'] == devices['JETSON']['proc'] and devices['JETSON']['found'])
                        ) else False
                        # only use mwt if the language supports it
                        stanza_lang = session['language_iso1']
                        stanza_has_mwt = False
                        try:
                            stanza_resources = stanza.resources.common.load_resources_json(os.getenv('STANZA_RESOURCES_DIR', stanza.resources.common.DEFAULT_MODEL_DIR))
                            stanza_has_mwt = 'mwt' in stanza_resources.get(stanza_lang, {})
                        except Exception:
                            pass
                        stanza_processors = 'tokenize,mwt,ner' if stanza_has_mwt else 'tokenize,ner'
                        stanza_nlp = stanza.Pipeline(stanza_lang, processors=stanza_processors, use_gpu=use_gpu, download_method=DownloadMethod.REUSE_RESOURCES, dir=os.getenv('STANZA_RESOURCES_DIR'))
                        if stanza_nlp:
                            session['stanza_cache'] = stanza_model
                            loaded_tts[stanza_model] = stanza_nlp
                            msg = f"NLP model {stanza_model} loaded!"
                            print(msg)
                except (ConnectionError, TimeoutError) as e:
                    error = f'Stanza model download connection error: {e}. Retry later'
                    print(error)
                    return []
                except Exception as e:
                    error = f'Stanza model initialization error: {e}'
                    print(error)
                    return []
            is_num2words_compat = get_num2words_compat(session['language_iso1'])
            try:
                with zipfile.ZipFile(session['epub_path'], 'r') as zf:
                    zip_names = set(zf.namelist())
                    zip_basenames = {os.path.basename(n): n for n in zip_names}
                    for doc_idx, doc in enumerate(all_docs):
                        text = filter_blocks(session_id, doc_idx, doc, stanza_nlp, is_num2words_compat, zf, zip_names, zip_basenames)
                        if text is None:
                            error = f'Error extracting content from document #{doc_idx + 1}; aborting conversion to avoid partial output.'
                            show_alert(session_id, {"type": "warning", "msg": error})
                            return []
                        blocks.append(text)
            finally:
                if stanza_nlp:
                    import gc, torch
                    try:
                        cache_key = session.get('stanza_cache')
                        if cache_key:
                            loaded_tts.pop(cache_key, None)
                        session['stanza_cache'] = None
                    except Exception:
                        pass
                    stanza_nlp = None
                    gc.collect()
                    if torch.cuda.is_available():
                        torch.cuda.empty_cache()
                        torch.cuda.ipc_collect()
            if len(blocks) == 0:
                error = 'No blocks found! possible reason: file corrupted or need to convert images to text with OCR'
                print(error)
                return []
            return blocks
        return []
    except Exception as e:
        error = f'Error extracting main content pages: {e}'
        DependencyError(error)
        return []

def filter_blocks(session_id:str, idx:int, doc:EpubHtml, stanza_nlp:Pipeline, is_num2words_compat:bool, zf:zipfile.ZipFile=None, zip_names:set=None, zip_basenames:dict=None)->str|None:
    def _tuple_row(node:Any, last_text_char:str|None=None)->Generator[tuple[str, Any], None, None]|None:
        try:
            prev_child_had_data = False
            for idx, child in enumerate(node.children):
                current_child_had_data = False
                if isinstance(child, NavigableString):
                    text = child.strip()
                    if text:
                        if prev_child_had_data:
                            yield ('break', sml_token("break"))
                        yield ('text', text)
                        last_text_char = text[-1]
                        current_child_had_data = True
                elif isinstance(child, Tag):
                    name = child.name.lower()
                    if name in heading_tags:
                        title = child.get_text(strip=True)
                        if title:
                            if prev_child_had_data:
                                yield ('break', sml_token("break"))
                            yield ('heading', title)
                            last_text_char = title[-1]
                            current_child_had_data = True
                    elif name == 'table':
                        if prev_child_had_data:
                            yield ('break', sml_token("break"))
                        yield ('table', child)
                        current_child_had_data = True
                    else:
                        return_data = False
                        if name in proc_tags:
                            is_header = False
                            if prev_child_had_data and name in break_tags:
                                yield ('break', sml_token("break"))
                            for inner in _tuple_row(child, last_text_char):
                                return_data = True
                                yield inner
                                if len(inner) > 1 and isinstance(inner[1], str) and inner[1]:
                                    last_text_char = inner[1][-1]
                                current_child_had_data = True
                                if inner[0] in ('text', 'heading') and isinstance(inner[1], str) and inner[1]:
                                    is_header = True
                            if return_data:
                                if name in break_tags and name != 'span':
                                    if is_header or (last_text_char and not last_text_char.isalnum() and not last_text_char.isspace()):
                                        yield ('break', sml_token("break"))
                                elif name in heading_tags or name in pause_tags:
                                    yield ('pause', sml_token("pause"))
                        else:
                            yield from _tuple_row(child, last_text_char)
                            current_child_had_data = True
                if current_child_had_data:
                    prev_child_had_data = True
        except Exception as e:
            error = f'filter_blocks() _tuple_row() error: {e}'
            DependencyError(error)
            return None

    def _num_repl(m):
        s = m.group(0)
        # leave years alone (already handled above)
        if re.fullmatch(r"\d{4}", s):
            return s
        n = float(s) if '.' in s else int(s)
        if is_num2words_compat:
            return num2words(n, lang=(lang_iso1 or 'en'))
        else:
            return math2words(m, lang, lang_iso1, tts_engine, is_num2words_compat)

    try:
        msg = f'----------\nParsing doc {idx}'
        print(msg)
        session = context.get_session(session_id)
        if session and session.get('id', False):
            lang, lang_iso1, tts_engine = session['language'], session['language_iso1'], session['tts_engine']
            heading_tags = [f'h{i}' for i in range(1, 5)]
            break_tags = ['br', 'p', 'span']
            pause_tags = ['div']
            proc_tags = heading_tags + break_tags + pause_tags
            doc_body = doc.get_body_content()
            raw_html = doc_body.decode('utf-8') if isinstance(doc_body, bytes) else doc_body
            soup = BeautifulSoup(raw_html, 'html.parser')
            body = soup.body
            if not body:
                msg = 'No body found. Skip to next doc…'
                print(msg)
                return ''
            # Skip known non-chapter types
            epub_type = body.get('epub:type', '').lower()
            if not epub_type:
                section_tag = soup.find('section')
                if section_tag:
                    epub_type = section_tag.get('epub:type', '').lower()
            excluded = {
                'frontmatter', 'backmatter', 'toc', 'titlepage', 'colophon',
                'acknowledgments', 'dedication', 'glossary', 'index',
                'appendix', 'bibliography', 'copyright-page', 'landmark'
            }
            if any(part in epub_type for part in excluded):
                msg = 'No body part. Skip to next doc…'
                print(msg)
                return ''
            # remove scripts/styles
            for tag in soup(['script', 'style']):
                tag.decompose()
            if not body.get_text(strip=True):
                images = body.find_all('img') + body.find_all('image')
                if images and zf:
                    msg = f'Doc {idx}: no text but {len(images)} image(s) detected. Running OCR…'
                    show_alert(session_id, {"type": "warning", "msg": msg})
                    ocr_parts = []
                    doc_dir = os.path.dirname(doc.get_name())
                    for img_tag in images:
                        img_ref = (
                            img_tag.get('src')
                            or img_tag.get('href')
                            or img_tag.get('{http://www.w3.org/1999/xlink}href')
                            or img_tag.get('xlink:href')
                        )
                        if not img_ref:
                            continue
                        img_zip_path = os.path.normpath(os.path.join(doc_dir, img_ref)).replace('\\', '/')
                        if img_zip_path not in zip_names:
                            img_zip_path = zip_basenames.get(os.path.basename(img_ref))
                        if not img_zip_path:
                            print(f'Could not resolve image in EPUB: {img_ref}')
                            continue
                        try:
                            img_data = zf.read(img_zip_path)
                            img = Image.open(io.BytesIO(img_data))
                            img = img.convert('RGB')
                            xhtml_content, error = ocr2xhtml(img, lang)
                            if xhtml_content:
                                ocr_parts.append(xhtml_content)
                            else:
                                show_alert(session_id, {"type": "warning", "msg": error})
                        except Exception as ocr_err:
                            print(f'OCR error on {img_zip_path}: {ocr_err}')
            tuples_list = list(_tuple_row(body))
            if not tuples_list:
                msg = 'No body text and no images found. Skip to next doc…'
                print(msg)
                return ''
            msg = f'Parsing xhtml markers…'
            print(msg)
            text_list = []
            handled_tables = set()
            prev_typ = None
            for typ, payload in tuples_list:
                if typ == 'heading':
                    text_list.append(payload.strip())
                elif typ in ('break', 'pause'):
                    if prev_typ != typ:
                        token = sml_token(typ)
                        if text_list and text_list[-1] not in {v['static'] for v in TTS_SML.values() if 'static' in v}:
                            text_list[-1] = text_list[-1] + token
                        else:
                            text_list.append(token)
                elif typ == 'table':
                    table = payload
                    if table in handled_tables:
                        prev_typ = typ
                        continue
                    handled_tables.add(table)
                    rows = table.find_all('tr')
                    if not rows:
                        prev_typ = typ
                        continue
                    headers = [c.get_text(strip=True) for c in rows[0].find_all(['td', 'th'])]
                    for row in rows[1:]:
                        cells = [c.get_text(strip=True).replace('\xa0', ' ') for c in row.find_all('td')]
                        if not cells:
                            continue
                        if len(cells) == len(headers) and headers:
                            line = ' — '.join(f'{h}: {c}' for h, c in zip(headers, cells))
                        else:
                            line = ' — '.join(cells)
                        if line:
                            text_list.append(line.strip())
                else:
                    text = payload.strip()
                    if text:
                        text_list.append(text)
                prev_typ = typ
            msg = f'Flattening as raw text…'
            print(msg)
            max_chars = int(language_mapping[lang]['max_chars'] / 1.5)
            clean_list = []
            i = 0
            while i < len(text_list):
                current = text_list[i]
                if current in {v['static'] for v in TTS_SML.values() if "static" in v}:
                    if clean_list:
                        prev = clean_list[-1]
                        if prev in {v['static'] for v in TTS_SML.values() if "static" in v}:
                            i += 1
                            continue
                    clean_list.append(current)
                    i += 1
                    continue
                clean_list.append(current)
                i += 1
            text = ' '.join(clean_list)
            if not re.search(r"[^\W_]", text):
                error = 'No valid text found!'
                print(error)
                return None
            # clean SML tags badly coded
            res, text = normalize_sml_tags(text)
            if res is False:
                show_alert(session_id, {"type": "warning", "msg": text})
                return None
            # remove any [break] between words or cutting words
            break_token = re.escape(sml_token('break'))
            strip_break_spaces_re = re.compile(rf'\s*{break_token}\s*')
            break_between_alnum_re = re.compile(rf'(?<=[\w]){break_token}(?=[\w])', flags=re.UNICODE)
            text = strip_break_spaces_re.sub(sml_token('break'), text)
            text = break_between_alnum_re.sub(' ', text)
            # escape all SML tags to not be touched by any text treatment
            text, sml_blocks = escape_sml(text)
            if stanza_nlp:
                msg = 'Converting dates and years to words…'
                print(msg)
                re_ordinal = re.compile(
                    r'(?<!\w)(0?[1-9]|[12][0-9]|3[01])(?:\s|\u00A0)*(?:st|nd|rd|th)(?!\w)',
                    re.IGNORECASE
                )
                re_num = re.compile(r'(?<!\w)[-+]?\d+(?:\.\d+)?(?!\w)')
                text = unicodedata.normalize('NFKC', text).replace('\u00A0', ' ')
                re_year = re.compile(r'\b(?:1[0-9]|20)\d{2}\b')
                if re_num.search(text) and (re_ordinal.search(text) or re_year.search(text)):
                    date_spans = get_date_entities(text, stanza_nlp)
                    if date_spans:
                        result = []
                        last_pos = 0
                        for start, end, date_text in date_spans:
                            result.append(text[last_pos:start])
                            # 1) convert 4-digit years (your original behavior)
                            processed = re.sub(
                                r"\b\d{4}\b",
                                lambda m: year2words(m.group(), lang, lang_iso1, is_num2words_compat),
                                date_text
                            )
                            # 2) convert ordinal days like "16th"/"16 th"->"sixteenth"
                            if is_num2words_compat:
                                processed = re_ordinal.sub(
                                    lambda m: num2words(int(m.group(1)), to='ordinal', lang=(lang_iso1 or 'en')),
                                    processed
                                )
                            else:
                                processed = re_ordinal.sub(
                                    lambda m: math2words(m.group(), lang, lang_iso1, tts_engine, is_num2words_compat),
                                    processed
                                )
                            # 3) convert other numbers (skip 4-digit years)
                            processed = re_num.sub(_num_repl, processed)
                            result.append(processed)
                            last_pos = end
                        result.append(text[last_pos:])
                        text = ' '.join(result)
                    else:
                        if is_num2words_compat:
                            text = re_ordinal.sub(
                                lambda m: num2words(int(m.group(1)), to='ordinal', lang=(lang_iso1 or 'en')),
                                text
                            )
                        else:
                            text = re_ordinal.sub(
                                lambda m: math2words(int(m.group(1)), lang, lang_iso1, tts_engine, is_num2words_compat),
                                text
                            )
                        text = re.sub(
                            r"\b\d{4}\b",
                            lambda m: year2words(m.group(), lang, lang_iso1, is_num2words_compat),
                            text
                        )
            msg = 'Convert romans to numbers…'
            print(msg)
            text = roman2number(text)
            msg = 'Convert time to words…'
            print(msg)
            text = clock2words(text, lang, lang_iso1, tts_engine, is_num2words_compat)
            msg = 'Convert numbers, maths signs to words…'
            print(msg)
            text = math2words(text, lang, lang_iso1, tts_engine, is_num2words_compat)
            msg = 'Normalize text…'
            print(msg)
            text = normalize_text(text, lang, lang_iso1, tts_engine)
            text = restore_sml(text, sml_blocks)
            return text
        return None
    except Exception as e:
        error = f'filter_blocks() error: {e}'
        DependencyError(error)
        return None

def get_sentences(session_id:str, text:str)->list|None:

    def split_inclusive(text:str, pattern:re.Pattern[str])->list[str]:
        result = []
        last_end = 0
        for match in pattern.finditer(text):
            result.append(text[last_end:match.end()].strip())
            last_end = match.end()
        if last_end < len(text):
            tail = text[last_end:].strip()
            if tail:
                result.append(tail)
        return result

    def split_sentence_on_sml(sentence:str)->list[str]:
        parts:list[str] = []
        last = 0
        for m in SML_TAG_PATTERN.finditer(sentence):
            start, end = m.span()
            if start > last:
                text = sentence[last:start]
                if text:
                    parts.append(text)
            parts.append(m.group(0))
            last = end
        if last < len(sentence):
            tail = sentence[last:]
            if tail:
                parts.append(tail)
        return parts

    def strip_escaped_sml(s:str)->str:
        return ''.join(c for c in s if ord(c) < sml_escape_tag)

    def clean_len(s:str)->int:
        return len(strip_escaped_sml(s))

    def is_latin_only(s:str)->bool:
        s = strip_escaped_sml(s)
        s = re.sub(r'[^\w\s]', '', s, flags=re.UNICODE)
        has_latin = bool(re.search(r'[A-Za-z]', s))
        has_nonlatin = bool(re.search(r'[^\x00-\x7F]', s))
        return has_latin and not has_nonlatin

    def split_at_space_limit(s:str)->list[str]:
        out = []
        rest = s.strip()
        while rest and len(strip_escaped_sml(rest)) > max_chars:
            cut = rest[:max_chars + 1]
            idx = cut.rfind(' ')
            if idx == -1:
                out.append(rest[:max_chars].strip())
                rest = rest[max_chars:].strip()
            else:
                out.append(rest[:idx].strip())
                rest = rest[idx + 1:].strip()
        if rest:
            out.append(rest.strip())
        return out

    def segment_ideogramms(text:str)->list[str]:
        result = []
        try:
            if lang in ['yue','yue-Hant','yue-Hans','zh-yue','cantonese']:
                import pycantonese as pc
                result.extend([t for t in pc.segment(text) if t.strip()])
            elif lang == 'zho':
                import jieba
                jieba.dt.cache_file = os.path.join(models_dir, 'jieba.cache')
                result.extend([t for t in jieba.cut(text) if t.strip()])
            elif lang == 'jpn':
                import nagisa
                result.extend(nagisa.tagging(text).words)
            elif lang == 'kor':
                from soynlp.tokenizer import LTokenizer
                ltokenizer = LTokenizer()
                result.extend([t for t in ltokenizer.tokenize(text) if t.strip()])
            elif lang in ['tha','lao','mya','khm']:
                from pythainlp.tokenize import word_tokenize
                result.extend([t for t in word_tokenize(text, engine='newmm') if t.strip()])
            else:
                result.append(text.strip())
            return result
        except Exception as e:
            DependencyError(e)
            return [text]

    def join_ideogramms(idg_list:list[str])->str:
        try:
            buffer = ''
            prev_latin = False
            prev_nonlatin = False
            for token in idg_list:
                cur_starts_latin = bool(re.match(r'[A-Za-z0-9]', token))
                cur_starts_nonlatin = bool(re.match(r'[^\x00-\x7F]', token))
                if buffer:
                    if (prev_latin and (cur_starts_latin or cur_starts_nonlatin)) or (prev_nonlatin and cur_starts_latin):
                        buffer += ' '
                    elif len(buffer) + len(token) > max_chars:
                        yield buffer
                        buffer = ''
                buffer += token
                prev_latin = bool(re.search(r'[A-Za-z0-9]$', token))
                prev_nonlatin = bool(re.search(r'[^\x00-\x7F]$', token))
            if buffer:
                yield buffer
        except Exception as e:
            DependencyError(e)
            if buffer:
                yield buffer

    try:
        session = context.get_session(session_id)
        if not session:
            return None

        lang, tts_engine = session['language'], session['tts_engine']
        max_chars = int(language_mapping[lang]['max_chars'] / 2)

        # escape all SML tags to not be touched by any text treatment
        text, sml_blocks = escape_sml(text)

        assert not SML_TAG_PATTERN.search(text)

        # PASS 1 — hard punctuation
        hard_pattern = re.compile(
            rf"(.*?(?:{'|'.join(map(re.escape, punctuation_split_hard_set))}))(?=\s|$)",
            re.DOTALL
        )
        hard_list = split_inclusive(text, hard_pattern)
        if not hard_list:
            hard_list = [text.strip()]
        hard_list = [s.strip() for s in hard_list if s.strip()]

        # PASS 2 — soft punctuation
        soft_pattern = re.compile(
            rf"(.*?(?:{'|'.join(map(re.escape, punctuation_split_soft_set))}))(?=\s|$)",
            re.DOTALL
        )
        soft_list = []
        i = 0
        n = len(hard_list)
        while i < n:
            s = hard_list[i].strip()
            if not s:
                i += 1
                continue
            if i + 1 < n:
                next_s = hard_list[i + 1].strip()
                next_clean = strip_escaped_sml(next_s)
                if next_clean and sum(c.isalnum() for c in next_clean) < 3:
                    s = f"{s} {next_s}"
                    i += 2
                else:
                    i += 1
            else:
                i += 1
            if len(strip_escaped_sml(s)) <= max_chars:
                soft_list.append(s)
                continue
            parts = split_inclusive(s, soft_pattern)
            if parts:
                valid = False
                for p in parts:
                    if len(strip_escaped_sml(p.strip())) <= max_chars:
                        valid = True
                        break
                if valid:
                    soft_list.extend([p.strip() for p in parts if p.strip()])
                else:
                    soft_list.append(s)
            else:
                soft_list.append(s)

        # PASS 3 — space split (last resort)
        last_list = []
        for s in soft_list:
            s = s.strip()
            if not s:
                continue
            rest = s
            while rest:
                current_len = len(strip_escaped_sml(rest))   # ← rename variable
                if current_len <= max_chars:
                    last_list.append(rest.strip())
                    break
                cut = rest[:max_chars + 1]
                idx = cut.rfind(' ')
                if idx > 0:
                    left = rest[:idx].strip()
                    right = rest[idx + 1:].strip()
                else:
                    left = rest[:max_chars].strip()
                    right = rest[max_chars:].strip()
                if not left or right == rest:
                    last_list.append(rest.strip())
                    break
                last_list.append(left)
                rest = right

        # PASS 4 — merge very short rows
        final_list = []
        merge_max_chars = int((max_chars / 2) / 3)
        i = 0
        n = len(last_list)
        while i < n:
            cur = last_list[i].strip()
            if not cur:
                i += 1
                continue
            if i == 0:
                final_list.append(cur)
                i += 1
                continue
            cur_len = clean_len(cur)
            if cur_len <= merge_max_chars:
                j = i + 1
                while j < n:
                    nxt = last_list[j].strip()
                    if not nxt:
                        j += 1
                        continue
                    if cur_len + clean_len(nxt) <= max_chars:
                        cur = cur.rstrip() + ' ' + nxt.lstrip()
                        cur_len = clean_len(cur)
                        j += 1
                        continue
                    break
                if final_list:
                    prev = final_list[-1]
                    if clean_len(prev) + cur_len <= max_chars:
                        final_list[-1] = prev.rstrip() + ' ' + cur.lstrip()
                        i = j
                        continue
                final_list.append(cur)
                i = j
                continue
            final_list.append(cur)
            i += 1

        if lang in ['zho', 'jpn', 'kor', 'tha', 'lao', 'mya', 'khm']:
            result = []
            for s in final_list:
                parts = split_sentence_on_sml(s)
                for part in parts:
                    part = part.strip()
                    if not part:
                        continue
                    if SML_TAG_PATTERN.fullmatch(part):
                        result.append(part)
                        continue
                    tokens = segment_ideogramms(part)
                    if isinstance(tokens, list):
                        result.extend([t for t in tokens if t.strip()])
                    else:
                        tokens = tokens.strip()
                        if tokens:
                            result.append(tokens)
            ideogram_list = []
            for s in join_ideogramms(result):
                if not is_latin_only(s):
                    ideogram_list.append(s)
            if ideogram_list:
                ideogram_list = [restore_sml(s, sml_blocks) for s in ideogram_list]
            return ideogram_list
        if final_list:
            final_list = [restore_sml(s, sml_blocks) for s in final_list]
        return final_list
    except Exception as e:
        print(f'get_sentences() error: {e}')
        return None

def get_sanitized(str:str, replacement:str='_')->str:
    str = str.replace('&', 'And')
    forbidden_chars = r'[<>:"/\\|?*\x00-\x1F ()]'
    sanitized = re.sub(r'\s+', replacement, str)
    sanitized = re.sub(forbidden_chars, replacement, sanitized)
    sanitized = sanitized.strip('_')
    return sanitized
    
def get_date_entities(text:str, stanza_nlp:Pipeline)->list[tuple[int,int,str]]|bool:
    try:
        doc = stanza_nlp(text)
        date_spans = []
        for ent in doc.ents:
            if ent.type == 'DATE':
                date_spans.append((ent.start_char, ent.end_char, ent.text))
        return date_spans
    except Exception as e:
        error = f'get_date_entities() error: {e}'
        print(error)
        return False

def get_num2words_compat(lang_iso1:str)->bool:
    try:
        test = num2words(1, lang=lang_iso1.replace('zh', 'zh_CN'))
        return True
    except NotImplementedError:
        return False
    except Exception as e:
        return False

def set_formatted_number(text:str, lang:str, lang_iso1:str, is_num2words_compat:bool, max_single_value:int=999_999_999_999_999_999)->str:
    # match up to 18 digits, optional “,…” groups (allowing spaces or NBSP after comma), optional decimal of up to 12 digits
    # handle optional range with dash/en dash/em dash between numbers, and allow trailing punctuation
    number_re = re.compile(
        r'(?<!\w)'
        r'(\d{1,18}(?:,\s*\d{1,18})*(?:\.\d{1,12})?)'      # first number
        r'(?:\s*([-–—])\s*'                                # dash type
        r'(\d{1,18}(?:,\s*\d{1,18})*(?:\.\d{1,12})?))?'    # optional second number
        r'([^\w\s]*)',                                     # optional trailing punctuation
        re.UNICODE
    )

    def normalize_commas(num_str:str)->str:
        # ormalize number string to standard comma format: 1,234,567
        tok = num_str.replace('\u00A0', '').replace(' ', '')
        if '.' in tok:
            integer_part, decimal_part = tok.split('.', 1)
            integer_part = integer_part.replace(',', '')
            integer_part = "{:,}".format(int(integer_part))
            return f'{integer_part}.{decimal_part}'
        else:
            integer_part = tok.replace(',', '')
            return "{:,}".format(int(integer_part))

    def clean_single_num(num_str:str)->str:
        tok = unicodedata.normalize('NFKC', num_str)
        if tok.lower() in ('inf', 'infinity', 'nan'):
            return tok
        clean = tok.replace(',', '').replace('\u00A0', '').replace(' ', '')
        try:
            num = float(clean) if '.' in clean else int(clean)
        except (ValueError, OverflowError):
            return tok
        if not math.isfinite(num) or abs(num) > max_single_value:
            return tok

        # Normalize commas before final output
        tok = normalize_commas(tok)

        if is_num2words_compat:
            new_lang_iso1 = lang_iso1.replace('zh', 'zh_CN')
            return num2words(num, lang=new_lang_iso1)
        else:
            phoneme_map = language_math_phonemes.get(
                lang,
                language_math_phonemes.get(default_language_code, language_math_phonemes['eng'])
            )
            return ' '.join(phoneme_map.get(ch, ch) for ch in str(num))

    def clean_match(match:re.Match)->str:
        first_num = clean_single_num(match.group(1))
        dash_char = match.group(2) or ''
        second_num = clean_single_num(match.group(3)) if match.group(3) else ''
        trailing = match.group(4) or ''
        if second_num:
            return f'{first_num}{dash_char}{second_num}{trailing}'
        else:
            return f'{first_num}{trailing}'

    return number_re.sub(clean_match, text)

def year2words(year_str:str, lang:str, lang_iso1:str, is_num2words_compat:bool)->str|bool:
    try:
        year = int(year_str)
        first_two = int(year_str[:2])
        last_two = int(year_str[2:])
        lang_iso1 = lang_iso1 if lang in language_math_phonemes.keys() else default_language_code
        lang_iso1 = lang_iso1.replace('zh', 'zh_CN')
        if not year_str.isdigit() or len(year_str) != 4 or last_two < 10:
            if is_num2words_compat:
                return num2words(year, lang=lang_iso1)
            else:
                return ' '.join(language_math_phonemes[lang].get(ch, ch) for ch in year_str)
        if is_num2words_compat:
            return f'{num2words(first_two, lang=lang_iso1)} {num2words(last_two, lang=lang_iso1)}' 
        else:
            return ' '.join(language_math_phonemes[lang].get(ch, ch) for ch in first_two) + ' ' + ' '.join(language_math_phonemes[lang].get(ch, ch) for ch in last_two)
    except Exception as e:
        error = f'year2words() error: {e}'
        print(error)
        return False

def clock2words(text:str, lang:str, lang_iso1:str, tts_engine:str, is_num2words_compat:bool)->str:

    def n2w(n:int)->str:
        key = (n, lang, is_num2words_compat)
        if key in _n2w_cache:
            return _n2w_cache[key]
        if is_num2words_compat:
            word = num2words(n, lang=lang_iso1)
        else:
            word = math2words(n, lang, lang_iso1, tts_engine, is_num2words_compat)
        if not isinstance(word, str):
            word = str(word)
        _n2w_cache[key] = word
        return word

    def repl_num(m:re.Match)->str:
        # Reject enumeration patterns like "(1.2)"
        start, end = m.start(), m.end()
        if (
            start > 0 and end < len(text)
            and text[start - 1] == '('
            and text[end] == ')'
        ):
            return m.group(0)
        # Parse hh[:mm[:ss]]
        try:
            h = int(m.group(1))
            mnt = int(m.group(2))
            sec = m.group(3)
            sec = int(sec) if sec is not None else None
        except Exception:
            return m.group(0)
        # basic validation; if out of range, keep original
        if not (0 <= h <= 23 and 0 <= mnt <= 59 and (sec is None or 0 <= sec <= 59)):
            return m.group(0)
        # If no language clock rules, just say numbers plainly
        if not lc:
            parts = [n2w(h)]
            if mnt != 0:
                parts.append(n2w(mnt))
            if sec is not None and sec > 0:
                parts.append(n2w(sec))
            return ' '.join(parts)
        next_hour = (h + 1) % 24
        special_hours = lc.get('special_hours', {})
        if mnt == 0 and (sec is None or sec == 0):
            if h in special_hours:
                phrase = special_hours[h]
            else:
                phrase = lc['oclock'].format(hour=n2w(h))
        elif mnt == 15:
            phrase = lc['quarter_past'].format(hour=n2w(h))
        elif mnt == 30:
            if lang == 'deu':
                phrase = lc['half_past'].format(next_hour=n2w(next_hour))
            else:
                phrase = lc['half_past'].format(hour=n2w(h))
        elif mnt == 45:
            phrase = lc['quarter_to'].format(next_hour=n2w(next_hour))
        elif mnt < 30:
            phrase = lc['past'].format(hour=n2w(h), minute=n2w(mnt)) if mnt != 0 else lc['oclock'].format(hour=n2w(h))
        else:
            minute_to_hour = 60 - mnt
            phrase = lc['to'].format(
                next_hour=n2w(next_hour),
                minute=n2w(minute_to_hour),
                minute_to_hour=n2w(minute_to_hour)
            )
        if sec is not None and sec > 0:
            second_phrase = lc['second'].format(second=n2w(sec))
            phrase = lc['full'].format(phrase=phrase, second_phrase=second_phrase)
        return phrase

    time_rx = re.compile(
        r'\b([01]?\d|2[0-3]):([0-5]\d)(?::([0-5]\d))?\b'
    )
    lc = language_clock.get(lang) if 'language_clock' in globals() else None
    _n2w_cache = {}
    return time_rx.sub(repl_num, text)

def math2words(text:str, lang:str, lang_iso1:str, tts_engine:str, is_num2words_compat:bool)->str:

    def repl_ambiguous(match:re.Match)->str:
        # handles "num SYMBOL num" and "SYMBOL num"
        if match.group(2) and match.group(2) in ambiguous_replacements:
            return f'{match.group(1)} {ambiguous_replacements[match.group(2)]} {match.group(3)}'
        if match.group(3) and match.group(3) in ambiguous_replacements:
            return f'{ambiguous_replacements[match.group(3)]} {match.group(4)}'
        return match.group(0)

    def _ordinal_to_words(m:re.Match)->str:
        n = int(m.group(1))
        if is_num2words_compat:
            try:
                return num2words(n, to='ordinal', lang=(lang_iso1 or 'en'))
            except Exception:
                pass
        # If num2words isn't available/compatible, keep original token as-is.
        return m.group(0)

    # Matches any digits + optional space/NBSP + st/nd/rd/th, not glued into words.
    re_ordinal = re.compile(r'(?<!\w)(\d+)(?:\s|\u00A0)*(?:st|nd|rd|th)(?!\w)')
    text = re.sub(r'(\d)\)', r'\1 : ', text)
    text = re_ordinal.sub(_ordinal_to_words, text)
    # Symbol phonemes
    ambiguous_symbols = {"-", "/", "*", "x"}
    phonemes_list = language_math_phonemes.get(lang, language_math_phonemes[default_language_code])
    replacements = {k: v for k, v in phonemes_list.items() if not k.isdigit() and k not in [',', '.']}
    normal_replacements  = {k: v for k, v in replacements.items() if k not in ambiguous_symbols}
    ambiguous_replacements = {k: v for k, v in replacements.items() if k in ambiguous_symbols}
    # Replace unambiguous symbols everywhere
    if normal_replacements:
        sym_pat = r'(' + '|'.join(map(re.escape, normal_replacements.keys())) + r')'
        text = re.sub(sym_pat, lambda m: f' {normal_replacements[m.group(1)]} ', text)
    # Replace ambiguous symbols only in valid equation contexts
    if ambiguous_replacements:
        ambiguous_pattern = (
            r'(?<!\S)'                   # no non-space before
            r'(\d+)\s*([-/*x])\s*(\d+)'  # num SYMBOL num
            r'(?!\S)'                    # no non-space after
            r'|'                         # or
            r'(?<!\S)([-/*x])\s*(\d+)(?!\S)'  # SYMBOL num
        )
        text = re.sub(ambiguous_pattern, repl_ambiguous, text)
    text = set_formatted_number(text, lang, lang_iso1, is_num2words_compat)
    return text

def roman2number(text:str)->str:

    def is_valid_roman(s:str)->bool:
        return bool(valid_roman.fullmatch(s))

    def to_int(s:str)->str:
        s = s.upper()
        i = 0
        result = 0
        while i < len(s):
            for roman, value in roman_numbers_tuples:
                if s[i:i + len(roman)] == roman:
                    result += value
                    i += len(roman)
                    break
            else:
                return s
        return str(result)

    def repl_heading(m: re.Match)->str:
        roman = m.group(1)
        if not is_valid_roman(roman):
            return m.group(0)
        return f"{to_int(roman)}{m.group(2)}{m.group(3)}"

    def repl_standalone(m: re.Match)->str:
        roman = m.group(1)
        if not is_valid_roman(roman):
            return m.group(0)
        return f"{to_int(roman)}{m.group(2)}"

    def repl_word(m: re.Match)->str:
        roman = m.group(1)
        if not is_valid_roman(roman):
            return m.group(0)
        return to_int(roman)

    def repl_chapter_single(m: re.Match)->str:
        word = m.group(1)
        roman = m.group(2)
        if not is_valid_roman(roman):
            return m.group(0)
        return f"{word} {to_int(roman)}"

    valid_roman = re.compile(
        r'^(?=.)M{0,3}(CM|CD|D?C{0,3})(XC|XL|L?X{0,3})(IX|IV|V?I{0,3})$',
        re.IGNORECASE
    )
    chapter_words = sorted(
        {w for words in chapter_word_mapping.values() for w in words},
        key=len,
        reverse=True
    )
    chapter_words_re = re.compile(
        rf'\b({"|".join(map(re.escape, chapter_words))})\s+([IVXLCDM]+)\b',
        re.IGNORECASE | re.UNICODE
    )
    text = re.sub(
        r'^(?:\s*)([IVXLCDM]+)([.-])(\s+)',
        repl_heading,
        text,
        flags=re.MULTILINE
    )
    text = re.sub(
        r'^(?:\s*)([IVXLCDM]+)([.-])(?:\s*)$',
        repl_standalone,
        text,
        flags=re.MULTILINE
    )
    text = chapter_words_re.sub(repl_chapter_single, text)
    text = re.sub(
        r'(?<!\S)([IVXLCDM]{2,})(?!\S)',
        repl_word,
        text
    )
    return text
    
def is_latin(s:str)->bool:
    return all((u'a' <= ch.lower() <= 'z') or ch.isdigit() or not ch.isalpha() for ch in s)

def foreign2latin(text:str, base_lang:str)->str:

    def script_of(word:str)->str:
        for ch in word:
            if ch.isalpha():
                name = unicodedata.name(ch, '')
                if 'CYRILLIC' in name:
                    return 'cyrillic'
                if 'LATIN' in name:
                    return 'latin'
                if 'ARABIC' in name:
                    return 'arabic'
                if 'HANGUL' in name:
                    return 'hangul'
                if 'HIRAGANA' in name or 'KATAKANA' in name:
                    return 'japanese'
                if 'CJK' in name or 'IDEOGRAPH' in name:
                    return 'chinese'
        return 'unknown'

    def romanize(word:str)->str:
        scr = script_of(word)
        if scr == 'latin':
            return word
        try:
            if scr == 'chinese':
                from pypinyin import pinyin, Style
                return ''.join(x[0] for x in pinyin(word, style=Style.NORMAL))
            if scr == 'japanese':
                import pykakasi
                k = pykakasi.kakasi()
                k.setMode('H', 'a')
                k.setMode('K', 'a')
                k.setMode('J', 'a')
                k.setMode('r', 'Hepburn')
                return k.getConverter().do(word)
            if scr == 'hangul':
                return unidecode(word)
            if scr == 'arabic':
                return unidecode(phonemize(word, language='ar', backend='espeak'))
            if scr == 'cyrillic':
                return unidecode(phonemize(word, language='ru', backend='espeak'))
            return unidecode(word)
        except Exception:
            return unidecode(word)

    # Protect ALL SML tags using the global grammar
    protected:dict[str, str] = {}
    for i, m in enumerate(SML_TAG_PATTERN.finditer(text)):
        key:str = f'__TTS_MARKER_{i}__'
        protected[key] = m.group(0)
        text = text.replace(m.group(0), key)
    tokens:list[str] = re.findall(r"\w+|[^\w\s]", text, re.UNICODE)
    buf:list[str] = []
    for t in tokens:
        if t in protected:
            buf.append(t)
        elif re.match(r"^\w+$", t):
            buf.append(romanize(t))
        else:
            buf.append(t)
    out:str = ''
    for i, t in enumerate(buf):
        if i == 0:
            out += t
        else:
            if re.match(r"^\w+$", buf[i - 1]) and re.match(r"^\w+$", t):
                out += ' ' + t
            else:
                out += t
    for k, v in protected.items():
        out = out.replace(k, v)
    return out

def normalize_sml_tags(text:str)->tuple[bool, str]:
    out = []
    stack = []
    last = 0
    for m in SML_TAG_PATTERN.finditer(text):
        start, end = m.span()
        out.append(text[last:start])
        tag = m.group("tag")
        close = bool(m.group("close"))
        value = m.group("value")
        info = TTS_SML.get(tag)
        if not info:
            out.append(m.group(0))
            last = end
            continue
        if info.get("paired"):
            if close:
                if not stack or stack[-1] != tag:
                    error = f'normalize_sml_tags() error: unmatched closing tag [/{tag}]'
                    return False, error
                stack.pop()
                out.append(f"[/{tag}]")
            else:
                stack.append(tag)
                if value is not None:
                    out.append(f"[{tag}:{value.strip()}]")
                else:
                    error = f'normalize_sml_tags() error: paired tag [{tag}] requires a value'
                    return False, error
        else:
            if close:
                error = f'normalize_sml_tags() error: non-paired tag [/{tag}] is invalid'
                return False, error
            out.append(info['static'])
        last = end
    out.append(text[last:])
    if stack:
        error = f"normalize_sml_tags() error: unclosed tag(s): {', '.join(stack)}"
        return False, error
    return True, ''.join(out)

def escape_sml(text:str)->tuple[str, list[str]]:
    sml_blocks:list[str] = []

    def replace(m:re.Match[str])->str:
        sml_blocks.append(m.group(0))
        return chr(sml_escape_tag + len(sml_blocks) - 1)

    return SML_TAG_PATTERN.sub(replace, text), sml_blocks

def restore_sml(text:str, sml_blocks:list[str])->str:
    for i, block in enumerate(sml_blocks):
        text = text.replace(chr(sml_escape_tag + i), block)
    return text

def sml_token(tag:str, value:str|None=None, close:bool=False)->str:
    if close:
        return f"[/{tag}]"
    if value is not None:
        return f"[{tag}:{value}]"
    return f"[{tag}]"

def normalize_text(text:str, lang:str, lang_iso1:str, tts_engine:str)->str:

    def replace(match:re.Match)->str:
        token = match.group(1)
        for k, expansion in mapping.items():
            if token.lower() == k.lower():
                return expansion
        return token  # fallback
            
    # Remove emojis
    emoji_pattern = re.compile(f"[{''.join(emojis_list)}]+", flags=re.UNICODE)
    emoji_pattern.sub('', text)
    if lang in abbreviations_mapping:
        mapping = abbreviations_mapping[lang]
        # Sort keys by descending length so longer ones match first
        keys = sorted(mapping.keys(), key=len, reverse=True)
        # Build a regex that only matches whole “words” (tokens) exactly
        pattern = re.compile(
            r'(?<!\w)(' + '|'.join(re.escape(k) for k in keys) + r')(?!\w)',
            flags=re.IGNORECASE
        )
        text = pattern.sub(replace, text)
    # This regex matches sequences like a., c.i.a., f.d.a., m.c., etc…
    pattern = re.compile(r'\b(?:[a-zA-Z]\.){1,}[a-zA-Z]?\b\.?')
    # uppercase acronyms
    text = re.sub(r'\b(?:[a-zA-Z]\.){1,}[a-zA-Z]?\b\.?', lambda m: m.group().replace('.', '').upper(), text)
    # romanize foreign words
    if language_mapping[lang]['script'] == 'latin':
        text = foreign2latin(text, lang)
    # Replace multiple newlines ("\n\n", "\r\r", "\n\r", etc.) with a [pause] 1.4sec
    pattern = r'(?:\r\n|\r|\n){2,}'
    text = re.sub(pattern, f" {sml_token('pause')} ", text)
    # Replace single newlines ("\n" or "\r") with spaces
    text = re.sub(r'\r\n|\r|\n', ' ', text)
    # Replace punctuations causing hallucinations
    pattern = f"[{''.join(map(re.escape, punctuation_switch.keys()))}]"
    text = re.sub(pattern, lambda match: punctuation_switch.get(match.group(), match.group()), text)
    # remove unwanted chars
    chars_remove_table = str.maketrans({ch: ' ' for ch in chars_remove})
    text = text.translate(chars_remove_table)
    # replace double quotes by a comma if no punctuation precedes it
    text = re.sub(r'\s*"\s*', '"', text)
    text = re.sub(r'(?<=[\p{L}\p{N}])"(?=[\p{L}\p{N}]|$)', ', ', text)
    text = re.sub(r'"', '', text)
    # Replace multiple and spaces with single space
    text = re.sub(r'\s+', ' ', text)
    # Replace ok by 'Owkey'
    text = re.sub(r'\bok\b', 'Okay', text, flags=re.IGNORECASE)
    # Escape special characters in the punctuation list for regex
    pattern = '|'.join(map(re.escape, punctuation_split_hard_set))
    # Reduce multiple consecutive punctuations hard
    text = re.sub(rf'(\s*({pattern})\s*)+', r'\2 ', text).strip()
    # Escape special characters in the punctuation list for regex
    pattern = '|'.join(map(re.escape, punctuation_split_soft_set))
    # Reduce multiple consecutive punctuations soft
    text = re.sub(rf'(\s*({pattern})\s*)+', r'\2 ', text).strip()
    # Pattern 1: Add a space between UTF-8 characters and numbers
    text = re.sub(r'(?<=[\p{L}])(?=\d)|(?<=\d)(?=[\p{L}])', ' ', text)
    # Replace special chars with words
    specialchars = specialchars_mapping.get(lang, specialchars_mapping.get(default_language_code, specialchars_mapping['eng']))
    specialchars_table = {ord(char): f" {word} " for char, word in specialchars.items()}
    text = text.translate(specialchars_table)
    text = ' '.join(text.split())
    return text

def block_hash(block: dict) -> str:
    return hashlib.sha1(
        '|'.join((
            (block.get('text') or '').strip(),
            block.get('voice') or '',
            block.get('tts_engine') or TTS_ENGINES['XTTSv2'],
            block.get('fine_tuned') or 'internal',
            json.dumps(block.get('sentences') or [], ensure_ascii=False),
        )).encode('utf-8')
    ).hexdigest()

def convert_chapters2audio(session_id:str)->bool:

    def _reset_chapter_file(block_id:str)->None:
        ch_file = os.path.join(session['chapters_dir'], f'{block_id}.{default_audio_proc_format}')
        if os.path.exists(ch_file):
            os.unlink(ch_file)
        block_dir = os.path.join(session['sentences_dir'], block_id)
        if os.path.isdir(block_dir):
            shutil.rmtree(block_dir)

    def _check_block_sentences(block_id:str, sentences:list)->set:
        block_dir = os.path.join(session['sentences_dir'], block_id)
        missing = set()
        for j in valid_idx:
            is_sml = bool(SML_TAG_PATTERN.fullmatch(sentences[j]))
            if (not is_sml) or (j == last_idx):
                sentence_file = os.path.join(block_dir, f'{j}.{default_audio_proc_format}')
                if not os.path.exists(sentence_file):
                    missing.add(j)
        return missing

    def _count_sentences(sentences:list)->int:
        return sum(1 for s in sentences if any(c.isalnum() for c in s.strip()))

    def _convert_ttsapi_sentence(sentence_file:str, sentence:str, block_voice:str|None)->tuple:
        manager = TTSManager(session)
        return manager.convert_sentence2audio(sentence_file, sentence, block_voice=block_voice)

    def _update_sentence_progress(sentence_text:str)->None:
        nonlocal global_sent
        global_sent += 1
        total_progress = (t.n + 1) / total_sentences
        if session['is_gui_process']:
            progress_bar(progress=total_progress, desc=f'{ebook_name} - {sentence_text}')
        t.set_description(f'{total_progress * 100:.2f}%')
        print(f' : {sentence_text}')
        t.update(1)

    session = context.get_session(session_id)
    if not (session and session.get('id', False)):
        return False
    try:
        if session['cancellation_requested']:
            return False
        print(f'*********** Session: {session_id} **************\n{session_info}')
        tts_manager = TTSManager(session)
        blocks_current = session['blocks_current']
        blocks = blocks_current['blocks']
        block_resume = blocks_current['block_resume']
        sentence_resume = blocks_current['sentence_resume']
        blocks_saved = session['blocks_saved']
        prev_blocks_list = blocks_saved.get('blocks', [])
        prev_blocks = {b['id']: b for b in prev_blocks_list} if isinstance(prev_blocks_list, list) else prev_blocks_list
        xtts_languages = default_engine_settings[TTS_ENGINES['XTTSv2']].get('languages', {})
        if session['language'] != 'eng' and session['language'] in xtts_languages:
            is_voice_changed = False
            voice_cache = {}
            for block in blocks:
                old_voice = block.get('voice')
                if old_voice in voice_cache:
                    new_voice = voice_cache[old_voice]
                else:
                    if old_voice is None:
                        new_voice = None
                    else:
                        new_voice, error = tts_manager.set_voice(old_voice)
                        if new_voice is None and error is not None:
                            show_alert(session_id, {'type': 'warning', 'msg': error})
                            return False
                    voice_cache[old_voice] = new_voice
                if new_voice != old_voice:
                    is_voice_changed = True
                    block['voice'] = new_voice
                    if blocks_saved:
                        if block['id'] in prev_blocks:
                            prev_blocks[block['id']]['voice'] = new_voice
            if is_voice_changed:
                if blocks_saved:
                    blocks_saved['blocks'] = list(prev_blocks.values())
                    session['blocks_saved'] = blocks_saved
                    save_json_blocks(session_id, 'blocks_saved')
                blocks_current['blocks'] = blocks
                session['blocks_current'] = blocks_current
                save_db_blocks(session_id)
        total_chapters = sum(1 for b in blocks if b['keep'] and b['text'].strip())
        if total_chapters == 0:
            show_alert(session_id, {'type': 'warning', 'msg': 'No chapters found!'})
            return False
        total_sentences = sum(_count_sentences(b['sentences']) for b in blocks if b['keep'] and b['text'].strip())
        if total_sentences == 0:
            show_alert(session_id, {'type': 'warning', 'msg': 'No sentences found!'})
            return False
        if not session['ebook']:
            return False
        ebook_name = Path(session['ebook']).name
        chapters_dir = session['chapters_dir']
        sentences_dir = session['sentences_dir']
        global_sent = 0
        ch_num = 0
        last_save_time = time.monotonic()
        baseline_initialized = False
        msg = (f'---------<br/>'
               f"{session['filename_noext']}<br/>"
               f"A total of {total_chapters} {'block' if total_chapters <= 1 else 'blocks'} "
               f"and {total_sentences} {'sentence' if total_sentences <= 1 else 'sentences'}."
               f'<br/>---------')
        show_alert(session_id, {'type': 'info', 'msg': msg})
        with tqdm(total=total_sentences, desc='0.00%', bar_format='{desc}: {n_fmt}/{total_fmt} ', unit='step', initial=0) as t:
            for x, block in enumerate(blocks):
                if not (block['keep'] and block['text'].strip()):
                    continue
                if session['cancellation_requested']:
                    return False
                ch_num += 1
                block_id = block['id']
                sentences = block['sentences']
                block_len = len(sentences)
                valid_idx = {i for i,s in enumerate(sentences) if any(c.isalnum() for c in s.strip())}
                last_idx = block_len - 1
                sent_start = global_sent
                current_hash = block_hash(block)
                block_ref = prev_blocks.get(block_id)
                hash_ref = block_hash(block_ref) if block_ref else None
                block_changed = block_ref is not None and hash_ref != current_hash
                missing_sentences = set()
                start_sentence = 0
                chapter_audio_file = os.path.join(chapters_dir, f'{block_id}.{default_audio_proc_format}')
                block_dir = os.path.join(sentences_dir, block_id)
                if x < block_resume and not block_changed:
                    if not os.path.exists(chapter_audio_file):
                        show_alert(session_id, {'type': 'warning', 'msg': f'Block {x} chapter audio missing, reconverting entire block…'})
                        _reset_chapter_file(block_id)
                    else:
                        missing_sentences = _check_block_sentences(block_id, sentences)
                        if not missing_sentences:
                            print(f'Chapter {ch_num} (block {x}) — has all sentences')
                            cnt = len(valid_idx)
                            global_sent += cnt
                            t.update(cnt)
                            continue
                        show_alert(session_id, {'type': 'warning', 'msg': f'Block {x} has {len(missing_sentences)} missing audio files, reconverting…'})
                        _reset_chapter_file(block_id)
                elif block_changed and x <= block_resume:
                    show_alert(session_id, {'type': 'info', 'msg': f'Chapter {ch_num} (block {x}) — changed, reconverting'})
                    _reset_chapter_file(block_id)
                elif x == block_resume and not block_changed:
                    if sentence_resume == 0 and os.path.isdir(block_dir):
                        shutil.rmtree(block_dir)
                    start_sentence = sentence_resume
                show_alert(session_id, {'type': 'info', 'msg': f'Chapter {ch_num} (block {x}) containing {block_len} sentences…'})
                os.makedirs(block_dir, exist_ok=True)
                blocks_current['block_resume'] = x
                blocks_current['sentence_resume'] = start_sentence
                session['blocks_current'] = blocks_current
                save_db_stamp(session_id)
                converted = False
                block_voice = block.get('voice') or session.get('voice')
                if session['tts_engine'] == TTS_ENGINES['TTSAPI']:
                    pending = []
                    pending_idx = set()
                    for j in range(block_len):
                        if j in valid_idx and (j >= start_sentence or j in missing_sentences):
                            sentence = sentences[j].strip()
                            sentence_file = os.path.join(block_dir, f'{j}.{default_audio_proc_format}')
                            pending.append((j, sentence, sentence_file))
                            pending_idx.add(j)
                    if pending and start_sentence > 0:
                        show_alert(session_id, {'type': 'info', 'msg': f'*** Resuming from sentence {global_sent} ***'})
                    for j in range(block_len):
                        if j in valid_idx and j not in pending_idx:
                            _update_sentence_progress(sentences[j].strip())
                    if pending:
                        if not baseline_initialized:
                            session['blocks_current'] = blocks_current
                            session['blocks_saved'] = copy.deepcopy(blocks_current)
                            save_json_blocks(session_id, 'blocks_saved')
                            baseline_initialized = True
                        workers = min(8, len(pending))
                        with ThreadPoolExecutor(max_workers=workers) as executor:
                            futures = {
                                executor.submit(_convert_ttsapi_sentence, sentence_file, sentence, block_voice): (j, sentence)
                                for j, sentence, sentence_file in pending
                            }
                            for future in as_completed(futures):
                                if session['cancellation_requested']:
                                    return False
                                j, sentence = futures[future]
                                run, error = future.result()
                                if not run:
                                    show_alert(session_id, {'type': 'warning', 'msg': error})
                                    return False
                                converted = True
                                _update_sentence_progress(sentence)
                        blocks_current['sentence_resume'] = block_len - 1
                        session['blocks_current'] = blocks_current
                        save_db_stamp(session_id)
                        last_save_time = time.monotonic()
                else:
                    for j in range(block_len):
                        if session['cancellation_requested']:
                            return False
                        sentence = sentences[j].strip()
                        if j in valid_idx:
                            if j >= start_sentence or j in missing_sentences:
                                if j == start_sentence and start_sentence > 0:
                                    show_alert(session_id, {'type': 'info', 'msg': f'*** Resuming from sentence {global_sent} ***'})
                                sentence_file = os.path.join(block_dir, f'{j}.{default_audio_proc_format}')
                                run, error = tts_manager.convert_sentence2audio(sentence_file, sentence, block_voice=block_voice)
                                if not run:
                                    show_alert(session_id, {'type': 'warning', 'msg': error})
                                    return False
                                converted = True
                                blocks_current['sentence_resume'] = j
                                now = time.monotonic()
                                if not baseline_initialized:
                                    session['blocks_current'] = blocks_current
                                    session['blocks_saved'] = copy.deepcopy(blocks_current)
                                    save_json_blocks(session_id, 'blocks_saved')
                                    baseline_initialized = True
                                elif now - last_save_time >= 5:
                                    session['blocks_current'] = blocks_current
                                    save_db_stamp(session_id)
                                    last_save_time = now
                            _update_sentence_progress(sentence)
                sent_end = global_sent - 1
                show_alert(session_id, {'type': 'info', 'msg': f'End of Chapter {ch_num} (block {x})'})
                if converted or block_changed or missing_sentences:
                    show_alert(session_id, {'type': 'info', 'msg': f'Combining chapter {ch_num} (block {x}) to audio, sentence {sent_start} to {sent_end}'})
                    session['blocks_current'] = blocks_current
                    save_db_stamp(session_id)
                    last_save_time = time.monotonic()
                    if not combine_audio_sentences(session_id, chapter_audio_file, block_id, block_len):
                        show_alert(session_id, {'type': 'warning', 'msg': 'combine_audio_sentences() failed!'})
                        return False
            blocks_current['block_resume'] = 0
            blocks_current['sentence_resume'] = 0
            session['blocks_current'] = blocks_current
            save_db_stamp(session_id)
            session['blocks_saved'] = copy.deepcopy(blocks_current)
            save_json_blocks(session_id, 'blocks_saved')
            return True
    except Exception as e:
        DependencyError(e)
        exception_alert(session_id, f'convert_chapters2audio() error: {e}')
        return False

def combine_audio_sentences(session_id:str, file:str, block_id:str, sentence_count:int)->bool:
    try:
        session = context.get_session(session_id)
        if not session or not session.get('id', False):
            error = 'Session expired!'
            print(error)
            return False
        if sentence_count == 0:
            error = f'No sentences to combine for block {block_id}.'
            print(error)
            return False
        block_dir = Path(session['sentences_dir']) / block_id
        ext = default_audio_proc_format
        selected_files = []
        missing = []
        for i in range(sentence_count):
            path = block_dir / f'{i}.{ext}'
            if path.is_file():
                selected_files.append(path)
            else:
                missing.append(i)
        if missing:
            error = f'Missing sentence files in block {block_id}: {missing}'
            print(error)
            return False
        concat_dir = session['process_dir']
        concat_list = os.path.join(concat_dir, 'concat_list_sentences.txt')
        with open(concat_list, 'w') as f:
            for path in selected_files:
                if session['cancellation_requested']:
                    return False
                f.write(f"file '{path.as_posix()}'\n")
        result = assemble_audio_chunks(concat_list, file, session['is_gui_process'])
        if not result:
            error = 'combine_audio_sentences() FFmpeg concat failed.'
            print(error)
            return False
        msg = f'********* Combined block audio file saved in {file}'
        print(msg)
        return True
    except Exception as e:
        DependencyError(e)
        return False

def combine_audio_chapters(session_id:str)->list[str]|None:

    def generate_ffmpeg_metadata(part_chapters:list[tuple[str,str]], output_metadata_path:str, default_audio_proc_format:str)->str|bool:
        try:
            out_fmt = session['output_format']
            is_mp4_like = out_fmt in ['mp4', 'm4a', 'm4b', 'mov']
            is_vorbis = out_fmt in ['ogg', 'webm']
            is_mp3 = out_fmt == 'mp3'
            def tag(key):
                return key.upper() if is_vorbis else key
            ffmpeg_metadata = ';FFMETADATA1\n'
            if session['metadata'].get('title'):
                ffmpeg_metadata += f"{tag('title')}={session['metadata']['title']}\n"
            if session['metadata'].get('creator'):
                ffmpeg_metadata += f"{tag('artist')}={session['metadata']['creator']}\n"
            if session['metadata'].get('language'):
                ffmpeg_metadata += f"{tag('language')}={session['metadata']['language']}\n"
            if session['metadata'].get('description'):
                ffmpeg_metadata += f"{tag('description')}={session['metadata']['description']}\n"
            if session['metadata'].get('publisher') and (is_mp4_like or is_mp3):
                ffmpeg_metadata += f"{tag('publisher')}={session['metadata']['publisher']}\n"
            if session['metadata'].get('published'):
                try:
                    if '.' in session['metadata']['published']:
                        year = datetime.strptime(session['metadata']['published'], '%Y-%m-%dT%H:%M:%S.%f%z').year
                    else:
                        year = datetime.strptime(session['metadata']['published'], '%Y-%m-%dT%H:%M:%S%z').year
                except Exception:
                    year = datetime.now().year
            else:
                year = datetime.now().year
            if is_vorbis:
                ffmpeg_metadata += f"{tag('date')}={year}\n"
            else:
                ffmpeg_metadata += f"{tag('year')}={year}\n"
            if session['metadata'].get('identifiers') and isinstance(session['metadata']['identifiers'], dict):
                if is_mp3 or is_mp4_like:
                    isbn = session['metadata']['identifiers'].get('isbn')
                    if isbn:
                        ffmpeg_metadata += f"{tag('isbn')}={isbn}\n"
                    asin = session['metadata']['identifiers'].get('mobi-asin')
                    if asin:
                        ffmpeg_metadata += f"{tag('asin')}={asin}\n"
            start_time = 0
            for filename, chapter_title in part_chapters:
                if session['cancellation_requested']:
                    return False
                filepath = os.path.join(session['chapters_dir'], filename)
                duration_ms = len(AudioSegment.from_file(filepath, format=default_audio_proc_format))
                clean_title = re.sub(r'(^#)|[=\\]|(-$)', lambda m: '\\' + (m.group(1) or m.group(0)), sanitize_meta_chapter_title(chapter_title))
                ffmpeg_metadata += '[CHAPTER]\nTIMEBASE=1/1000\n'
                ffmpeg_metadata += f'START={start_time}\nEND={start_time + duration_ms}\n'
                ffmpeg_metadata += f"{tag('title')}={clean_title}\n"
                start_time += duration_ms
            with open(output_metadata_path, 'w', encoding='utf-8') as f:
                f.write(ffmpeg_metadata)
            return output_metadata_path
        except Exception as e:
            error = f'generate_ffmpeg_metadata() Error: {e}'
            print(error)
            return False

    def export_audio(combined_audio:str, metadata_file:str, final_file:str, block_indices:set=None, part_num:int=None)->bool:

        def on_progress(p:float)->None:
            if is_gui_process:
                progress_bar(p / 100.0, desc=f'Export Part {part_num}' if part_num is not None else 'Export')

        is_gui_process = session['is_gui_process']
        try:
            if session['cancellation_requested']:
                return False
            ffprobe_cmd = [
                shutil.which('ffprobe'), '-v', 'error', '-threads', '0', '-select_streams', 'a:0',
                '-show_entries', 'stream=codec_name,sample_rate,sample_fmt',
                '-of', 'default=nokey=1:noprint_wrappers=1', combined_audio
            ]
            probe = subprocess.run(ffprobe_cmd, capture_output=True, text=True)
            if probe.returncode != 0:
                error = f'ffprobe failed for {combined_audio}: {probe.stderr.strip()}'
                print(error)
                return False
            codec_info = probe.stdout.strip().splitlines()
            input_codec = codec_info[0] if len(codec_info) > 0 else None
            input_rate = codec_info[1] if len(codec_info) > 1 else None
            cmd = [shutil.which('ffmpeg'), '-hide_banner', '-nostats', '-hwaccel', 'auto', '-thread_queue_size', '1024', '-i', combined_audio]
            target_codec, target_rate = None, None
            if session['output_format'] == 'wav':
                target_codec = 'pcm_s16le'
                target_rate = '44100'
                cmd += ['-map', '0:a', '-ar', target_rate, '-sample_fmt', 's16']
            elif session['output_format'] == 'aac':
                target_codec = 'aac'
                target_rate = '44100'
                cmd += ['-c:a', 'aac', '-b:a', '192k', '-ar', target_rate, '-movflags', '+faststart']
            elif session['output_format'] == 'flac':
                target_codec = 'flac'
                target_rate = '44100'
                cmd += ['-c:a', 'flac', '-compression_level', '5', '-ar', target_rate]
            else:
                cmd += ['-f', 'ffmetadata', '-i', metadata_file, '-map', '0:a']
                if session['output_format'] in ['m4a', 'm4b', 'mp4', 'mov']:
                    target_codec = 'aac'
                    target_rate = '44100'
                    cmd += ['-c:a', 'aac', '-b:a', '192k', '-ar', target_rate, '-movflags', '+faststart+use_metadata_tags']
                elif session['output_format'] == 'mp3':
                    target_codec = 'mp3'
                    target_rate = '44100'
                    cmd += ['-c:a', 'libmp3lame', '-b:a', '192k', '-ar', target_rate]
                elif session['output_format'] == 'webm':
                    target_codec = 'opus'
                    target_rate = '48000'
                    cmd += ['-c:a', 'libopus', '-b:a', '192k', '-ar', target_rate]
                elif session['output_format'] == 'ogg':
                    target_codec = 'opus'
                    target_rate = '48000'
                    cmd += ['-c:a', 'libopus', '-compression_level', '0', '-b:a', '192k', '-ar', target_rate]
                cmd += ['-map_metadata', '1']
            if session['output_channel'] == 'stereo':
                cmd += ['-ac', '2']
            else:
                cmd += ['-ac', '1']
            if input_codec == target_codec and input_rate == target_rate:
                cmd = [
                    shutil.which('ffmpeg'), '-hide_banner', '-nostats', '-hwaccel', 'auto', '-thread_queue_size', '1024', '-i', combined_audio,
                    '-threads', '0', '-f', 'ffmetadata', '-i', metadata_file,
                    '-map', '0:a', '-map_metadata', '1', '-c', 'copy',
                    '-progress', 'pipe:2',
                    '-y', final_file
                ]
            else:
                cmd += [
                    '-filter_threads', '0',
                    '-filter_complex_threads', '0',
                    '-af', 'loudnorm=I=-16:LRA=11:TP=-1.5:linear=true,afftdn=nf=-70',
                    '-threads', '0',
                    '-progress', 'pipe:2',
                    '-y', final_file
                ]
            proc_pipe = SubprocessPipe(cmd, is_gui_process=is_gui_process, total_duration=get_audio_duration(combined_audio), msg='Export', on_progress=on_progress)
            if not proc_pipe.result:
                error = f'ffmpeg export failed for {final_file}'
                print(error)
                return False
            if not (os.path.exists(final_file) and os.path.getsize(final_file) > 0):
                error = f'{Path(final_file).name} is corrupted or does not exist'
                print(error)
                return False
            if session['output_format'] in ['mp3', 'm4a', 'm4b', 'mp4'] and session['cover'] is not None:
                cover_path = session['cover']
                msg = f'Adding cover {cover_path} into the final audiobook file…'
                print(msg)
                audio = None
                if session['output_format'] == 'mp3':
                    from mutagen.mp3 import MP3
                    from mutagen.id3 import ID3, APIC, error as id3_error
                    audio = MP3(final_file, ID3=ID3)
                    try:
                        audio.add_tags()
                    except id3_error:
                        pass
                    with open(cover_path, 'rb') as img:
                        audio.tags.add(APIC(encoding=3, mime='image/jpeg', type=3, desc='Cover', data=img.read()))
                elif session['output_format'] in ['mp4', 'm4a', 'm4b']:
                    from mutagen.mp4 import MP4, MP4Cover
                    audio = MP4(final_file)
                    with open(cover_path, 'rb') as f:
                        cover_data = f.read()
                    audio['covr'] = [MP4Cover(cover_data, imageformat=MP4Cover.FORMAT_JPEG)]
                if audio is not None:
                    audio.save()
            final_vtt = os.path.join(session['audiobooks_dir'], f'{Path(final_file).stem}.vtt')
            vtt_built, error = build_vtt_file(session, vtt_path=final_vtt, block_indices=block_indices)
            if not vtt_built:
                error = f'build_vtt_file() error: {error}'
                print(error)
                return False
            return True
        except Exception as e:
            error = f'Export failed: {e}'
            print(error)
            return False

    try:
        session = context.get_session(session_id)
        if not (session and session.get('id', False)):
            return None
        chapter_files = []
        chapter_titles = []
        chapter_positions = []
        for x, block in enumerate(session['blocks_current']['blocks']):
            if not (block['keep'] and block['text'].strip()):
                continue
            if not block.get('sentences'):
                error = f"Block {x} (id {block['id']}) has no sentences but is marked keep"
                print(error)
                return None
            block_id = block['id']
            fname = f'{block_id}.{default_audio_proc_format}'
            fpath = os.path.join(session['chapters_dir'], fname)
            if not os.path.exists(fpath):
                error = f'Missing chapter audio for block {x} (id {block_id}): {fpath}'
                print(error)
                return None
            chapter_files.append(fname)
            chapter_titles.append(block['sentences'][0])
            chapter_positions.append(x)
        is_gui_process = session['is_gui_process']
        if len(chapter_files) == 0:
            print('No block files exist!')
            return None
        chunks_size = 892
        total_duration = 0.0
        durations = []
        for i in range(0, len(chapter_files), chunks_size):
            filepaths = [
                os.path.join(session['chapters_dir'], f)
                for f in chapter_files[i:i + chunks_size]
            ]
            durations_dict = get_audiolist_duration(filepaths)
            for path in filepaths:
                dur = durations_dict.get(path, 0.0)
                durations.append(dur)
                total_duration += dur
        if len(durations) != len(chapter_files):
            error = f'Duration count mismatch: {len(durations)} durations vs {len(chapter_files)} chapter files'
            print(error)
            return None
        exported_files = []
        concat_dir = session['process_dir']
        if session.get('output_split'):
            part_files = []
            part_chapter_indices = []
            cur_part = []
            cur_indices = []
            cur_duration = 0
            max_part_duration = int(session['output_split_hours']) * 3600
            for idx, (file, dur) in enumerate(zip(chapter_files, durations)):
                if session['cancellation_requested']:
                    return None
                if cur_part and (cur_duration + dur > max_part_duration):
                    part_files.append(cur_part)
                    part_chapter_indices.append(cur_indices)
                    cur_part = []
                    cur_indices = []
                    cur_duration = 0
                cur_part.append(file)
                cur_indices.append(idx)
                cur_duration += dur
            if cur_part:
                part_files.append(cur_part)
                part_chapter_indices.append(cur_indices)
            pad_width = len(str(len(part_files)))
            is_multi_part = len(part_files) > 1
            for part_idx, (part_file_list, indices) in enumerate(zip(part_files, part_chapter_indices)):
                concat_list = os.path.join(concat_dir, f'concat_list_chapters_{part_idx+1:0{pad_width}d}.txt')
                with open(concat_list, 'w') as f:
                    for file in part_file_list:
                        if session['cancellation_requested']:
                            return None
                        path = Path(session['chapters_dir']) / file
                        f.write(f"file '{path.as_posix()}'\n")
                merged_audio = Path(session['process_dir']) / f"{get_sanitized(session['metadata']['title'])}_part{part_idx+1:0{pad_width}d}.{default_audio_proc_format}"
                result = assemble_audio_chunks(concat_list, merged_audio, is_gui_process)
                if not result:
                    error = f'assemble_audio_chunks() Final merge failed for part {part_idx+1}.'
                    print(error)
                    return None
                metadata_file = Path(session['process_dir']) / f'metadata_part{part_idx+1:0{pad_width}d}.txt'
                part_chapters = [(chapter_files[i], chapter_titles[i]) for i in indices]
                generate_ffmpeg_metadata(part_chapters, str(metadata_file), default_audio_proc_format)
                final_file = os.path.join(
                    session['audiobooks_dir'],
                    f"{Path(session['final_name']).stem}_part{part_idx+1:0{pad_width}d}.{session['output_format']}"
                    if is_multi_part else session['final_name']
                )
                block_indices = {chapter_positions[i] for i in indices} if is_multi_part else None
                if export_audio(merged_audio, metadata_file, final_file, block_indices=block_indices, part_num=part_idx+1):
                    exported_files.append(final_file)
        else:
            concat_list = os.path.join(concat_dir, 'concat_list_chapters_1.txt')
            merged_audio = Path(session['process_dir']) / f"{get_sanitized(session['metadata']['title'])}.{default_audio_proc_format}"
            with open(concat_list, 'w') as f:
                for file in chapter_files:
                    if session['cancellation_requested']:
                        return None
                    path = Path(session['chapters_dir']) / file
                    f.write(f"file '{path.as_posix()}'\n")
            result = assemble_audio_chunks(concat_list, merged_audio, is_gui_process)
            if not result:
                print(f'assemble_audio_chunks() Final merge failed for {merged_audio}.')
                return None
            metadata_file = os.path.join(session['process_dir'], 'metadata.txt')
            chapters_zip = list(zip(chapter_files, chapter_titles))
            generate_ffmpeg_metadata(chapters_zip, metadata_file, default_audio_proc_format)
            final_file = os.path.join(session['audiobooks_dir'], session['final_name'])
            if export_audio(merged_audio, metadata_file, final_file):
                exported_files.append(final_file)
        return exported_files if exported_files else None
    except Exception as e:
        DependencyError(e)
        return None

def assemble_audio_chunks(txt_file:str, out_file:str, is_gui_process:bool)->bool:

    def on_progress(p:float)->None:
        if is_gui_process:
            progress_bar(p / 100.0, desc='Assemble')

    try:
        total_duration = 0.0
        filepaths = []
        try:
            with open(txt_file, 'r') as f:
                for line in f:
                    if line.strip().startswith('file'):
                        file_path = (
                            line.strip()
                            .split('file ')[1]
                            .strip()
                            .strip("'")
                            .strip('"')
                        )
                        if os.path.exists(file_path):
                            filepaths.append(file_path)
            durations = get_audiolist_duration(filepaths)
            total_duration = sum(durations.values())
        except Exception as e:
            error = f'assemble_audio_chunks() open file {txt_file} Error: {e}'
            print(error)
            return False
        ffmpeg = shutil.which('ffmpeg')
        if not ffmpeg:
            error = 'ffmpeg not found'
            print(error)
            return False
        cmd = [
            ffmpeg,
            '-hide_banner',
            '-nostats',
            '-safe', '0',
            '-f', 'concat',
            '-i', txt_file,
            '-c:a', default_audio_proc_format,
            '-map_metadata', '-1',
            '-threads', '0',
            '-progress', 'pipe:2',
            '-y', out_file
        ]
        proc_pipe = SubprocessPipe(
            cmd=cmd,
            is_gui_process=is_gui_process,
            total_duration=total_duration,
            msg='Assemble',
            on_progress=on_progress
        )
        if proc_pipe.result and os.path.exists(out_file):
            msg = f'Completed → {out_file}'
            print(msg)
            return True
        else:
            error = f'Failed (proc_pipe) → {out_file}'
            print(error)
            return False
    except subprocess.CalledProcessError as e:
        DependencyError(e)
        return False
    except Exception as e:
        error = f'assemble_audio_chunks() Error: Failed to process {txt_file} → {out_file}: {e}'
        print(error)
        return False

def ellipsize_utf8_bytes(s:str, max_bytes:int, ellipsis:str='…')->str:
    s = '' if s is None else str(s)
    if max_bytes <= 0:
        return ''
    raw = s.encode('utf-8')
    e = ellipsis.encode('utf-8')
    if len(raw) <= max_bytes:
        return s
    if len(e) >= max_bytes:
        # return as many bytes of the ellipsis as fit
        return e[:max_bytes].decode('utf-8', errors='ignore')
    budget = max_bytes - len(e)
    out = bytearray()
    for ch in s:
        b = ch.encode('utf-8')
        if len(out) + len(b) > budget:
            break
        out.extend(b)
    return out.decode('utf-8') + ellipsis

def sanitize_meta_chapter_title(title:str, max_bytes:int=140)->str:
    # avoid None and embedded NULs which some muxers accidentally keep
    title = (title or '').replace('\x00', '')
    title = title.replace(sml_token('pause'), '')
    return ellipsize_utf8_bytes(title, max_bytes=max_bytes, ellipsis='…')

def delete_proc_audio_files(dir:str, files:list)->None:
    base = Path(dir)
    for file in base.rglob(f'[0-9]*.{default_audio_proc_format}'):
        if file.stem.isdigit() and file in files:
            file.unlink()

def delete_folder(folder_path:str)->None:
    for name in os.listdir(folder_path):
        path = os.path.join(folder_path, name)
        if os.path.isfile(path) or os.path.islink(path):
            os.unlink(path)
        else:
            shutil.rmtree(path)

def delete_unused_tmp_dirs(session_id:str, output_dir:str, days:int)->None:
    session = context.get_session(session_id)
    if session and session.get('id', False):
        dir_array = [
            tmp_dir,
            output_dir,
            os.path.join(models_dir, '__sessions'),
            os.path.join(voices_dir, '__sessions')
        ]
        current_user_dirs = {
            f'proc-{session_id}',
            f'web-{session_id}',
            f'cli-{session_id}',
            f'voice-{session_id}',
            f'model-{session_id}'
        }
        current_time = time.time()
        threshold_time = current_time - (days * 24 * 60 * 60)  # Convert days to seconds
        for dir_path in dir_array:
            if os.path.exists(dir_path) and os.path.isdir(dir_path):
                for dir in os.listdir(dir_path):
                    if dir in current_user_dirs:        
                        full_dir_path = os.path.join(dir_path, dir)
                        if os.path.isdir(full_dir_path):
                            try:
                                dir_mtime = os.path.getmtime(full_dir_path)
                                dir_ctime = os.path.getctime(full_dir_path)
                                if dir_mtime < threshold_time and dir_ctime < threshold_time:
                                    shutil.rmtree(full_dir_path, ignore_errors=True)
                                    msg = f'Deleted expired session: {full_dir_path}'
                                    print(msg)
                            except Exception as e:
                                error = f'Error deleting {full_dir_path}: {e}'
                                print(error)

def get_compatible_tts_engines(language:str)->list[str]:
    return [
        engine
        for engine, cfg in default_engine_settings.items()
        if language in cfg.get('languages', {})
    ]

def is_ttsapi_engine(engine:str|None)->bool:
    return engine == TTS_ENGINES.get('TTSAPI')

def is_voice_file_path(engine:str|None, voice:Any)->bool:
    return isinstance(voice, str) and os.path.exists(voice) and not is_ttsapi_engine(engine)

def resolve_voice(session_id:str, ebook_src:str)->str|None:
    """
    Returns the voice to use for a given ebook in DIRECTORY mode.
    Lookup order: voice_map[abs(ebook_src)] -> voice_map[basename] -> session['voice'] -> None.
    Voice file existence is checked; missing overrides fall through to the default.
    """
    session = context.get_session(session_id)
    if not session:
        return None
    if not ebook_src:
        return session.get('voice')
    voice_map = session.get('voice_map') or {}
    abs_src = os.path.abspath(ebook_src)
    override = voice_map.get(abs_src) or voice_map.get(os.path.basename(ebook_src))
    if is_ttsapi_engine(session.get('tts_engine')):
        return override if override not in (None, '') else session.get('voice')
    if override and os.path.exists(override):
        return override
    return session.get('voice')

def convert_ebook(args:dict)->tuple:
    try:
        global context
        error = None
        session_id = None
        info_session = None
        if not args.get('id'):
            error = 'Session ID is missing!'
            return error, False
        session_id = str(args['id'])
        session = context.get_session(session_id)
        if not session or (session and not session.get('id', False)):
            error = 'Session expired or does not exist!'
            return error, False
        if args['language'] is not None:
            try:
                if len(args['language']) in (2, 3):
                    lang_dict = Lang(args['language'])
                    if lang_dict:
                        args['language'] = lang_dict.pt3
                        args['language_iso1'] = lang_dict.pt1
                else:
                    args['language_iso1'] = None
            except Exception as e:
                pass
            if args['language'] not in language_mapping.keys():
                error = 'The language you provided is not (yet) supported'
                return error, False
            if args.get('ebook_mode') == ebook_modes['TEXT']:
                if not args['ebook_textarea']:
                    error = 'Ebook textarea is empty.'
                    return error, False
                text = args['ebook_textarea']
                text_name = get_sanitized(text[:64])
                text_name_hash = hashlib.md5(f'{text_name}_{session_id}'.encode()).hexdigest()
                text_filename = SML_TAG_PATTERN.sub('', text_name)
                text_filename = get_sanitized(text_filename)
                text_filename = f'{text_filename[:48]}_{session_id}.txt'
                text_filepath = os.path.join(tempfile.gettempdir(), text_filename)
                with open(text_filepath, 'w', encoding='utf-8') as f:
                    f.write(text)
                session['ebook_textarea'] = args['ebook_textarea']
                session['ebook_src'] = text_filepath
            else:
                if args.get('ebook_src'):
                    if not os.path.splitext(args['ebook_src'])[1]:
                        error = f"{args['ebook_src']} needs a format extension."
                        return error, False
                    if not os.path.exists(args['ebook_src']):
                        error = 'File does not exist or Directory empty.'
                        return error, False
                    session['ebook_src'] = str(args['ebook_src'])
            session['custom_model_dir'] = os.path.join(models_dir, '__sessions',f"model-{session_id}")
            session['script_mode'] = str(args['script_mode']) if args.get('script_mode') is not None else NATIVE
            session['is_gui_process'] = bool(args['is_gui_process'])
            session['blocks_preview'] = bool(args['blocks_preview']) if args.get('blocks_preview') else False
            session['device'] = str(args['device'])
            session['language'] = str(args['language'])
            session['language_iso1'] = str(args['language_iso1'])
            session['tts_engine'] = str(args['tts_engine']) if args.get('tts_engine') is not None else str(get_compatible_tts_engines(args['language'])[0])
            session['custom_model'] =  args['custom_model']
            session['fine_tuned'] = str(args['fine_tuned'])
            session['voice'] = args.get('voice', None)
            session['xtts_temperature'] =  float(args['xtts_temperature'])
            session['xtts_length_penalty'] = float(args['xtts_length_penalty'])
            session['xtts_num_beams'] = int(args['xtts_num_beams'])
            session['xtts_repetition_penalty'] = float(args['xtts_repetition_penalty'])
            session['xtts_top_k'] =  int(args['xtts_top_k'])
            session['xtts_top_p'] = float(args['xtts_top_p'])
            session['xtts_speed'] = float(args['xtts_speed'])
            session['xtts_enable_text_splitting'] = bool(args['xtts_enable_text_splitting'])
            session['bark_text_temp'] =  float(args['bark_text_temp'])
            session['bark_waveform_temp'] =  float(args['bark_waveform_temp'])
            session['output_format'] = str(args['output_format'])
            session['output_channel'] = str(args['output_channel'])
            session['output_split'] = bool(args['output_split'])
            session['output_split_hours'] = args['output_split_hours']if args['output_split_hours'] is not None else default_output_split_hours
            session['model_cache'] = f"{session['tts_engine']}-{session['fine_tuned']}"
            session['session_dir'] = os.path.join(tmp_dir, f'proc-{session_id}')
            session['status'] = status_tags['EDIT'] if session['blocks_preview'] else status_tags['CONVERTING'] 
            ebook_name = get_sanitized(Path(session['ebook_src']).stem)
            cleanup_models_cache()
            print(f"Processing eBook file: {os.path.basename(session['ebook_src'])}")
            if session['is_gui_process']:
                session['final_name'] = ebook_name + '.' + session['output_format']
                session['process_dir'] = os.path.join(session['session_dir'], f"{hashlib.md5(os.path.join(session['audiobooks_dir'], Path(session['final_name']).stem).encode()).hexdigest()}")
                session['chapters_dir'] = os.path.join(session['process_dir'], "chapters")
                session['sentences_dir'] = os.path.join(session['chapters_dir'], 'sentences')
            else:
                session['system'] = DEVICE_SYSTEM
                session['audiobooks_dir'] = os.path.abspath(args['output_dir']) if args.get('output_dir') is not None else os.path.join(audiobooks_cli_dir, f'cli-{session_id}')
                session['final_name'] = os.path.join(session['audiobooks_dir'], ebook_name + '.' + session['output_format'])
                session['process_dir'] = os.path.join(session['session_dir'], f"{hashlib.md5(os.path.join(session['audiobooks_dir'], Path(session['final_name']).stem).encode()).hexdigest()}")
                session['chapters_dir'] = os.path.join(session['process_dir'], "chapters")
                session['sentences_dir'] = os.path.join(session['chapters_dir'], 'sentences')
                session['voice_dir'] = os.path.join(voices_dir, '__sessions', f'voice-{session_id}', session['language'])
                os.makedirs(session['voice_dir'], exist_ok=True)
                audio_pre_final_exist = os.path.exists(os.path.join(session['process_dir'], ebook_name + '.' + default_audio_proc_format))
                audio_sentences_exist = any(Path(session['sentences_dir']).rglob(f'*.{default_audio_proc_format}'))
                if audio_pre_final_exist or audio_sentences_exist:
                    msg = f"Warning! This conversion already exists. Continue? WARNING! The whole previous conversion will be deleted!" if audio_pre_final_exist else f"Warning! Some sentences are already converted. Resume?"
                    print(msg)
                    while True:
                        choice = input("[s]kip / [y]es: ").strip().lower()
                        if choice in ('s', 'y'):
                            break
                        print("Please enter 's', or 'y'.")
                    if choice == 'y':
                        if audio_pre_final_exist:
                            delete_folder(session['process_dir'])
                    elif choice == 's':
                        msg = 'Conversion skipped.'
                        return msg, True
                if error is None:
                    #delete_unused_tmp_dirs(session_id, audiobooks_cli_dir, tmp_expire)
                    if session['custom_model'] is not None:
                        if not os.path.exists(session['custom_model_dir']):
                            os.makedirs(session['custom_model_dir'], exist_ok=True)
                        custom_src_path = Path(session['custom_model'])
                        custom_src_name = custom_src_path.stem
                        if not os.path.exists(os.path.join(session['custom_model_dir'], custom_src_name)):
                            try:
                                if analyze_uploaded_file(session['custom_model'], default_engine_settings[session['tts_engine']]['files']):
                                    model = extract_custom_model(session_id)
                                    if model is not None:
                                        session['custom_model'] = model
                                    else:
                                        error = f"{model} could not be extracted or mandatory files are missing"
                                else:
                                    error = f'{os.path.basename(f)} is not a valid model or some required files are missing'
                            except ModuleNotFoundError as e:
                                error = f"No presets module for TTS engine '{session['tts_engine']}': {e}"
                    if session.get('voice'):
                        if is_ttsapi_engine(session['tts_engine']):
                            if os.path.exists(str(session['voice'])):
                                error = 'TTSAPI voice override must be a model id, not a local voice file path.'
                        else:
                            voice_name = os.path.splitext(os.path.basename(session['voice']))[0].replace('&', 'And')
                            voice_name = get_sanitized(voice_name)
                            final_voice_file = os.path.join(session['voice_dir'], f'{voice_name}.wav')
                            if not os.path.exists(final_voice_file):
                                extractor = VoiceExtractor(session, session['voice'], voice_name)
                                voice_status, msg = extractor.extract_voice()
                                if voice_status:
                                    session['voice'] = final_voice_file
                                else:
                                    error = f'VoiceExtractor.extract_voice() failed! {msg}'
            if error is None:
                if session['script_mode'] == NATIVE:
                    is_installed = check_programs('Calibre', 'ebook-convert', '--version')
                    if not is_installed:
                        error = f'check_programs() Calibre failed: {e}'
                    is_installed = check_programs('FFmpeg', 'ffmpeg', '-version')
                    if not is_installed:
                        error = f'check_programs() FFMPEG failed: {e}'
                if error is None:
                    if prepare_dirs(session_id):
                        session['ebook'] = os.path.join(session['process_dir'], os.path.basename(session['ebook_src']))
                        shutil.copy(session['ebook_src'], session['ebook'])
                        session['filename_noext'] = os.path.splitext(os.path.basename(session['ebook']))[0]
                        msg = ''
                        msg_extra = ''                      
                        if session['device'] == devices['CUDA']['proc']:
                            if not devices['CUDA']['found']:
                                session['device'] = devices['CPU']['proc']
                                msg += f'CUDA not supported by the Torch installed!<br/>Read {default_gpu_wiki}<br/>Switching to CPU'
                        elif session['device'] == devices['JETSON']['proc'] or session['device'] == devices['JETSON']['proc']:
                            if not devices['JETSON']['found']:
                                session['device'] = devices['CPU']['proc']
                                msg += f'JETSON CUDA not supported by the Torch installed!<br/>Read {default_gpu_wiki}<br/>Switching to CPU'
                        elif session['device'] == devices['MPS']['proc']:
                            if not devices['MPS']['found']:
                                session['device'] = devices['CPU']['proc']
                                msg += f'MPS not supported by the Torch installed!<br/>Read {default_gpu_wiki}<br/>Switching to CPU'
                        elif session['device'] == devices['ROCM']['proc']:
                            if not devices['ROCM']['found']:
                                session['device'] = devices['CPU']['proc']
                                msg += f'ROCM not supported by the Torch installed!<br/>Read {default_gpu_wiki}<br/>Switching to CPU'
                        elif session['device'] == devices['XPU']['proc']:
                            if not devices['XPU']['found']:
                                session['device'] = devices['CPU']['proc']
                                msg += f"XPU not supported by the Torch installed!<br/>Read {default_gpu_wiki}<br/>Switching to CPU"
                        vram_dict = VRAMDetector().detect_vram(session['device'], session['script_mode'])
                        print(f'vram_dict: {vram_dict}')
                        total_vram_gb = vram_dict.get('total_vram_gb', 0)
                        detected_free_vram_gb = vram_dict.get('free_vram_gb', 0)
                        session['free_vram_gb'] = detected_free_vram_gb
                        if session['free_vram_gb'] == 0:
                            msg_extra += f"<br/>Memory capacity not detected! restrict to {session['free_vram_gb']}GB max"
                        else:
                            msg_extra += f"<br/>Free Memory available: {session['free_vram_gb']}GB"
                            if session['free_vram_gb'] < default_engine_settings[session['tts_engine']]['rating']['VRAM']:
                                msg_extra += f"<br/>Free Memory {session['free_vram_gb']} is lower than VRAM/RAM {default_engine_settings[session['tts_engine']]['rating']['VRAM']}GB required!<br/>It will probably crash the conversion!"
                            if session['free_vram_gb'] > 4.0:
                                if session['tts_engine'] == TTS_ENGINES['BARK']:
                                    os.environ['SUNO_USE_SMALL_MODELS'] = 'False'  
                        if session['tts_engine'] == TTS_ENGINES['BARK']:
                            if session['free_vram_gb'] < 12.0:
                                os.environ['SUNO_OFFLOAD_CPU'] = "True"
                                os.environ['SUNO_USE_SMALL_MODELS'] = "True"
                                msg_extra += f"<br/>Switching BARK to SMALL models"  
                            else:
                                os.environ['SUNO_OFFLOAD_CPU'] = "False"
                                os.environ['SUNO_USE_SMALL_MODELS'] = "False"
                        if msg == '':
                            msg_extra = f"Using {session['device'].upper()}" + msg_extra
                        device_vram_required = default_engine_settings[session['tts_engine']]['rating']['RAM'] if session['device'] == devices['CPU']['proc'] else default_engine_settings[session['tts_engine']]['rating']['VRAM']
                        if float(total_vram_gb) >= float(device_vram_required):
                            if msg:
                                show_alert(session_id, {"type": "warning", "msg": msg + msg_extra})
                            else:
                                show_alert(session_id, {"type": "info", "msg": msg_extra})
                            session['epub_path'] = os.path.join(session['process_dir'], f"__{session['filename_noext']}.epub")
                            session['blocks_orig_json'] = os.path.join(session['process_dir'], f"{file_prefixes['clone']}{session['filename_noext']}.json")
                            session['blocks_saved_json']   = os.path.join(session['process_dir'], f"{file_prefixes['saved']}{session['filename_noext']}.json")
                            session['blocks_current_db']   = os.path.join(session['process_dir'], f"{file_prefixes['current']}{session['filename_noext']}.db")
                            checksum, error = compare_checksums(session_id)
                            if not checksum or not os.path.exists(session['epub_path']):
                                result_epub = convert2epub(session_id)
                                if result_epub:
                                    if os.path.exists(session['epub_path']):
                                        for jf in (session['blocks_orig_json'], session['blocks_saved_json']):
                                            if os.path.exists(jf):
                                                os.unlink(jf)
                                        db = session['blocks_current_db']
                                        for f in (db, db + '-wal', db + '-shm'):
                                            if os.path.exists(f):
                                                os.unlink(f)
                                        msg = f"NOTE: process folder {session['process_dir']} is strictly used for internal tasks and has nothing to do with the final conversion."
                                        print(msg)
                                    else:
                                        error = f"convert2epub() {session['epub_path']} does not exists! check write permissions."
                                else:
                                    error = 'convert2epub() error: could not convert to epub file!'
                            if error is None:
                                missing_orig_json = True
                                if os.path.exists(session['blocks_orig_json']):
                                    missing_orig_json = False
                                    blocks_orig = load_json_blocks(session['blocks_orig_json'])
                                    is_changed = False
                                    is_reset = False
                                    if blocks_orig:
                                        blocks = blocks_orig.get('blocks', [])
                                        new_blocks = []
                                        for block in blocks:
                                            if any(c.isalnum() for c in block.get('text','')):
                                                if not block.get('id'):
                                                    block['id'] = str(uuid.uuid4())
                                                    is_changed = True
                                                new_blocks.append(block)
                                            else:
                                                is_reset = True
                                        blocks_orig['blocks'] = new_blocks
                                        session['blocks_orig'] = blocks_orig
                                    if is_changed or is_reset:
                                        save_json_blocks(session_id, 'blocks_orig')
                                    if os.path.exists(session['blocks_saved_json']):
                                        blocks_saved = load_json_blocks(session['blocks_saved_json'])
                                        if blocks_saved:
                                            session['blocks_saved'] = blocks_saved
                                            if is_changed or is_reset:
                                                if is_changed:
                                                    blocks = blocks_saved.get('blocks', [])
                                                    for i, block in enumerate(blocks):
                                                        if i < len(blocks_orig['blocks']):
                                                            block['id'] = blocks_orig['blocks'][i]['id']
                                                    blocks_saved['blocks'] = blocks
                                                    session['blocks_saved'] = blocks_saved
                                                elif is_reset:
                                                    session['blocks_saved'] = copy.deepcopy(blocks_orig)
                                                save_json_blocks(session_id, 'blocks_saved')
                                    if os.path.exists(session['blocks_current_db']):
                                        blocks_current = load_db_blocks(session['blocks_current_db'])
                                        if blocks_current:
                                            session['blocks_current'] = blocks_current
                                            if is_changed or is_reset:
                                                if is_changed:
                                                    blocks = blocks_current.get('blocks', [])
                                                    for i, block in enumerate(blocks):
                                                        if i < len(blocks_orig['blocks']):
                                                            block['id'] = blocks_orig['blocks'][i]['id']
                                                    blocks_current['blocks'] = blocks
                                                    session['blocks_current'] = blocks_current
                                                elif is_reset:
                                                    session['blocks_current'] = copy.deepcopy(blocks_orig)
                                                save_db_blocks(session_id)
                                epubBook = epub.read_epub(session['epub_path'], {'ignore_ncx': True})
                                if epubBook:
                                    metadata = dict(session['metadata'])
                                    for key, value in metadata.items():
                                        data = epubBook.get_metadata('DC', key)
                                        if data:
                                            for value, attributes in data:
                                                metadata[key] = value
                                    metadata['language'] = session['language']
                                    metadata['title'] = metadata['title'] or Path(session['ebook']).stem.replace('_', ' ')
                                    metadata['creator'] = False if not metadata['creator'] or metadata['creator'] == 'Unknown' else metadata['creator']
                                    session['metadata'] = metadata
                                    try:
                                        if len(session['metadata']['language']) == 2:
                                            lang_dict = Lang(session['language'])
                                            if lang_dict:
                                                session['metadata']['language'] = lang_dict.pt3
                                    except Exception as e:
                                        pass
                                    if session['metadata']['language'] != session['language']:
                                        error = f"WARNING!!! language selected {session['language']} differs from the EPUB file language {session['metadata']['language']}"
                                        show_alert(session_id, {'type': 'warning', 'msg': error})
                                    is_lang_in_tts_engine = (
                                        session.get('tts_engine') in default_engine_settings and
                                        session.get('language') in default_engine_settings[session['tts_engine']].get('languages', {})
                                    )
                                    if is_lang_in_tts_engine:
                                        session['cover'] = get_cover(epubBook, session_id)
                                        if session.get('cover', False):
                                            if missing_orig_json:
                                                raw_blocks = get_blocks(session_id, epubBook)
                                                if raw_blocks:
                                                    session['blocks_orig'] = {
                                                        "page": 0,
                                                        "block_resume": 0,
                                                        "sentence_resume": 0,
                                                        "voice": session['voice'],
                                                        "tts_engine": session['tts_engine'],
                                                        "fine_tuned": session['fine_tuned'],
                                                        "blocks": [
                                                            {
                                                                "id": str(uuid.uuid4()),
                                                                "expand": False,
                                                                "keep": True,
                                                                "text": t,
                                                                "voice": session['voice'],
                                                                "tts_engine": session['tts_engine'],
                                                                "fine_tuned": session['fine_tuned'],
                                                                "sentences": [],
                                                            }
                                                            for t in raw_blocks if t
                                                        ],
                                                    }
                                                if session.get('blocks_orig', {}):
                                                    save_json_blocks(session_id, 'blocks_orig')
                                            if not session.get('blocks_current', {}):
                                                session['blocks_current'] = copy.deepcopy(session['blocks_orig'])
                                                save_db_blocks(session_id)
                                            # --- legacy upgrade: old snapshots may lack top-level scalars (TO REMOVE AFTER A WHILE) ---
                                            for key in ('blocks_orig', 'blocks_current', 'blocks_saved'):
                                                snap = session.get(key)
                                                if snap:
                                                    changed = False
                                                    if 'voice' not in snap:
                                                        snap['voice'] = session.get('voice')
                                                        snap['tts_engine'] = session.get('tts_engine')
                                                        snap['fine_tuned'] = session.get('fine_tuned')
                                                        changed = True
                                                    if 'page' not in snap:
                                                        snap['page'] = 0
                                                        changed = True
                                                    if changed:
                                                        session[key] = snap
                                                        if key == 'blocks_current':
                                                            save_db_blocks(session_id)
                                                        else:
                                                            save_json_blocks(session_id, key)
                                            # --------------------------------#
                                            if session.get('blocks_orig', {}) and session.get('blocks_current', {}):
                                                sync_globals_to_blocks(session_id)
                                                if session['blocks_preview']:
                                                    msg = f'Chapters preview requested. Select which block to convert:'
                                                    print(msg)
                                                    progress_status = os.path.basename(session['ebook'])
                                                    return progress_status, True
                                                else:
                                                    progress_status, passed = finalize_audiobook(session_id)
                                                    return progress_status, passed
                                            else:
                                                error = f"get_blocks() or save_json_blocks() failed! {session['blocks_orig']}"
                                        else:
                                            error = 'get_cover() failed!'
                                    else:
                                        error = f"language {session['language']} not supported by {session['tts_engine']}!"
                                else:
                                    error = 'epubBook.read_epub failed!'
                        else:
                            error = f"Your device has not enough memory ({total_vram_gb}GB) to run {session['tts_engine']} engine ({device_vram_required}GB)"
                    else:
                        error = f"Temporary directory {session['process_dir']} not removed due to failure."
        if session['cancellation_requested']:
            error = 'Conversion Cancelled'
        return error, False
    except Exception as e:
        error = f'convert_ebook() Exception: {e}'
        return error, False

def finalize_audiobook(session_id:str)->tuple:
    try:
        session = context.get_session(session_id)
        is_preview = session.get('blocks_preview', False) if session else False
        result = lambda msg, ok: (gr.update(value=msg), gr.update(value=ok)) if is_preview else (msg, ok)

        def fail(error):
            session['status'] = status_tags['END']
            return result(error, False)

        if not session or not session.get('id', False):
            msg = 'session expired!'
            return result(msg, False)
        if session['status'] not in [status_tags['EDIT'], status_tags['CONVERTING']]:
            msg = 'No blocks have been selected for the conversion!'
            return result(msg, False)
        if not session.get('blocks_current', {}):
            error = 'finalize_audiobook() failed! blocks_current empty!'
            return fail(error)
        session['status'] = status_tags['CONVERTING']
        print('Get sentences…')
        blocks_current = session['blocks_current']
        blocks = blocks_current['blocks']
        for idx, block in enumerate(blocks):
            if session['cancellation_requested']:
                if session['status'] == status_tags['DISCONNECTED']:
                    context_tracker.end_session(session_id, session['socket_hash'])
                    msg = 'Frontend disconnected!'
                    return result(msg, False)
                msg = 'Conversion cancelled'
                return result(msg, False)
            if not block['keep'] or not block['text'].strip():
                block['sentences'] = []
                continue
            if block.get('sentences', []):
                print(f'Block {idx} — sentences already split, skipping')
                continue
            sentences_list = get_sentences(session_id, block['text'])
            if sentences_list is None:
                error = 'No sentences found!'
                return result(error, False)
            block['sentences'] = sentences_list
        blocks_current['blocks'] = blocks
        session['blocks_current'] = blocks_current
        conversion = convert_chapters2audio(session_id)
        if not conversion:
            error = 'convert_chapters2audio() failed!'
            session = context.get_session(session_id)
            if session and session.get('id', False):
                if session['cancellation_requested']:
                    error = 'Conversion cancelled'
            return fail(error)
        show_alert(session_id, {'type': 'info', 'msg': 'Combining sentences and chapters…'})
        exported_files = combine_audio_chapters(session_id)
        if exported_files is None:
            return fail('combine_audio_chapters() error: exported_files not created!')
        session['audiobook'] = exported_files[-1]
        filename = os.path.basename(session['ebook'])
        count_ebook = 0
        if session['ebook_mode'] == ebook_modes['DIRECTORY']:
            if isinstance(session['ebook_list'], list):
                if session['ebook_src'] in session['ebook_list']:
                    ebook_list = session['ebook_list']
                    ebook_list.remove(session['ebook_src'])
                    session['ebook_list'] = ebook_list
                count_ebook = len(session['ebook_list'])
        if count_ebook > 0:
            reset_ebook_session(session_id, force=True, filter_keys=False)
            show_alert(session_id, {'type': 'success', 'msg': f'{filename} / converted. {count_ebook} ebook(s) conversion remaining…'})
        else:
            if session['ebook_mode'] == ebook_modes['DIRECTORY']:
                session['ebook_list'] = None
                session['voice_map'] = {}
                session['ebook_selected'] = None
            elif session['ebook_mode'] == ebook_modes['SINGLE']:
                session['ebook_src'] = None
            elif session['ebook_mode'] == ebook_modes['TEXT']:
                ebook_src = session['ebook_src']
                if ebook_src:
                    try:
                        os.remove(ebook_src)
                    except FileNotFoundError:
                        pass
                    except OSError:
                        pass
                session['ebook_src'] = session['ebook_src_notextarea']
            session['status'] = status_tags['END']
            reset_ebook_session(session_id, force=True, filter_keys=False)
            show_alert(session_id, {'type': 'success', 'msg': f'{filename} / converted.'})
            print(f'*********** Session: {session_id} **************\n{session_info}')
        return result(filename, True)
    except Exception as e:
        session['status'] = status_tags['END']
        reset_ebook_session(session_id, force=True, filter_keys=False)
        DependencyError(e)
        error = f'finalize_audiobook(): {e}'
        exception_alert(session_id, error)
        return result(error, False)

def on_unload(req:gr.Request)->None:
    socket_hash = req.session_hash
    if any(socket_hash in session for session in context.sessions.values()):
        session_id = context.find_id_by_hash(socket_hash)
        if session_id:
            session = context.get_session(session_id)
            if session['status'] == status_tags['CONVERTING']:
                session['cancellation_requested'] = True
                session['status'] = status_tags['DISCONNECTED']
                session['socket_hash'] = socket_hash
            else:
                context_tracker.end_session(session_id, socket_hash)

def restore_session_from_data(data:dict, session:DictProxy, force:bool, filter_keys:bool)->None:
    try:
        for key, value in data.items():
            if key in session:
                if filter_keys and key in save_session_keys_except:
                    continue
                if isinstance(value, dict) and value and isinstance(session[key], dict):
                    nested = session[key]
                    restore_session_from_data(value, nested, force, filter_keys)
                    session[key] = nested
                else:
                    if not force:
                        if value is None and session[key] is not None:
                            continue
                    session[key] = value
    except Exception as e:
        DependencyError(e)
        error = f'restore_session_from_data(): {e}'
        exception_alert(session['id'], error)

def reset_ebook_session(session_id:str, force:bool, filter_keys:bool)->None:
    session = context.get_session(session_id)
    data = {
        "ebook": None,
        "process_dir": None,
        "chapters_dir": None,
        "sentences_dir": None,
        "epub_path": None,
        "final_name": None,
        "filename_noext": None,
        "cover": None,
        "blocks_orig": {},
        "blocks_saved": {},
        "blocks_current": {},
        "blocks_orig_json": None,
        "blocks_saved_json": None,
        "blocks_current_db": None,
        "audiobook_overridden": None,
        "metadata": {
            "title": None, 
            "creator": None,
            "contributor": None,
            "language": None,
            "identifier": None,
            "publisher": None,
            "date": None,
            "description": None,
            "subject": None,
            "rights": None,
            "format": None,
            "type": None,
            "coverage": None,
            "relation": None,
            "Source": None,
            "Modified": None,
        }
    }
    restore_session_from_data(data, session, force, filter_keys=filter_keys)

def cleanup_models_cache()->None:
    try:
        active_models = {
            cache
            for session in context.sessions.values()
            for cache in (session.get('model_cache'), session.get('model_zs_cache'), session.get('stanza_cache'))
            if cache is not None
        }
        for key in list(loaded_tts.keys()):
            if key not in active_models:
                del loaded_tts[key]
        gc.collect()
    except Exception as e:
        error = f"cleanup_models_cache() error: {e}"
        print(error)

def show_alert(session_id:str|None, state:dict|None)->None:
    if state is not None:
        if state.get('msg'):
            print(state['msg'].replace('<br/>', '\n'))
            if session_id is not None:
                session = context.get_session(session_id)
                if session.get('is_gui_process'):
                    if isinstance(state, dict):
                        if state['type'] is not None:
                            if state['type'] == 'error':
                                raise gr.Error(state['msg'])
                            elif state['type'] == 'warning':
                                gr.Warning(state['msg'])
                            elif state['type'] == 'info':
                                gr.Info(state['msg'])
                            elif state['type'] == 'success':
                                gr.Success(state['msg'])

def exception_alert(session_id:str|None, error:str|None)->None:
    if error is not None:
        print(error.replace('<br/>', '\n'))
        if session_id is not None:
            session = context.get_session(session_id)
            if session and session.get('id', False):
                session['status'] = status_tags['READY']
                if session['is_gui_process']:
                    raise gr.Error(error)

def get_all_ip_addresses()->list:
    ip_addresses = []
    for interface, addresses in psutil.net_if_addrs().items():
        for address in addresses:
            if address.family in [socket.AF_INET, socket.AF_INET6]:
                ip_addresses.append(address.address)
    return ip_addresses
